"""Streamlit interface for the broiler chicken financial model."""

from __future__ import annotations

import copy
from dataclasses import asdict, dataclass, is_dataclass
from datetime import datetime, timezone
from html import escape as xml_escape
import hashlib
from io import BytesIO
import json
import math
import re
import shutil
import subprocess
import tempfile
from xml.etree import ElementTree as ET
import zipfile
from typing import Any, Dict, Iterable, List, Optional, Tuple, get_type_hints

import pandas as pd
import streamlit as st
import altair as alt
from streamlit.delta_generator import DeltaGenerator
from streamlit.errors import StreamlitAPIException

from broiler_model.assumptions import Assumptions, ASSUMPTION_SCHEDULE_LAYOUT
from broiler_model.analytics import (
    AnalyticsPlan,
    break_even_analysis,
    build_predictive_analytics,
    goal_seek_live_price,
    perform_what_if_analysis,
    run_custom_simulations,
    run_monte_carlo_analysis,
    scenario_planning,
)
from broiler_model.config import (
    load_custom_simulation_definitions,
    load_monte_carlo_distributions,
)
from broiler_model.model import generate_model_outputs
from broiler_model.production import summarise_revenue_totals


DEFAULT_CUSTOM_SIMULATION_DEFINITIONS = load_custom_simulation_definitions()
DEFAULT_MONTE_CARLO_DISTRIBUTIONS = load_monte_carlo_distributions()

RAG_DOCUMENT_TYPE_CONFIG: Dict[str, Dict[str, Any]] = {
    "pdf": {
        "extensions": {"pdf"},
        "pipeline": ["pypdf/PyPDF2", "pdfplumber", "PyMuPDF", "pdfminer", "pdftotext", "OCR"],
    },
    "word": {
        "extensions": {"docx"},
        "pipeline": ["python-docx", "docx-zipxml"],
    },
    "excel": {
        "extensions": {"xlsx", "xlsm", "xltx", "xltm", "xls"},
        "pipeline": ["openpyxl/pandas", "xlsx-zipxml"],
    },
    "delimited": {
        "extensions": {"csv", "tsv"},
        "pipeline": ["utf-8 delimited parser"],
    },
    "json": {
        "extensions": {"json"},
        "pipeline": ["json parser"],
    },
    "markup": {
        "extensions": {"html", "htm", "xml"},
        "pipeline": ["tag-strip text extraction"],
    },
    "pptx": {
        "extensions": {"pptx"},
        "pipeline": ["python-pptx", "pptx-zipxml"],
    },
    "text": {
        "extensions": {"txt", "md", "log"},
        "pipeline": ["plain text"],
    },
    "archive": {
        "extensions": {"zip"},
        "pipeline": ["zip expansion into supported files"],
    },
}


ROW_REMOVAL_COLUMN = "Remove row"
ROW_EDIT_COLUMN = "Edit row"


_ASSUMPTION_LABEL_TO_KEY = {
    label: key for _, label, key in ASSUMPTION_SCHEDULE_LAYOUT
}
try:
    _ASSUMPTION_FIELD_TYPES = get_type_hints(Assumptions)
except Exception:  # pragma: no cover - defensive fallback
    _ASSUMPTION_FIELD_TYPES = {
        name: field.type for name, field in Assumptions.__dataclass_fields__.items()
    }


def _to_numeric(series: pd.Series) -> pd.Series:
    """Coerce a pandas Series to numeric values where possible."""

    return pd.to_numeric(series, errors="coerce")


def _to_float(value: Any) -> Optional[float]:
    """Return float when value is numeric-like, otherwise ``None``."""

    if value is None:
        return None
    if isinstance(value, (int, float)):
        return float(value)
    if isinstance(value, str):
        stripped = value.strip()
        if not stripped:
            return None
        try:
            return float(stripped)
        except ValueError:
            return None
    try:
        return float(value)
    except (TypeError, ValueError):
        return None




def _decode_bytes_to_text(data: bytes) -> str:
    """Decode byte payloads with lightweight encoding fallbacks."""

    for encoding in ("utf-8", "utf-8-sig", "utf-16", "utf-16-le", "utf-16-be", "latin-1"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="ignore")


def _tokenize_text(text: str) -> List[str]:
    """Tokenize text into lowercase alphanumeric terms."""

    return [token for token in re.findall(r"[a-z0-9]+", (text or "").lower()) if token]


def _rag_supported_extensions(include_archives: bool = True) -> List[str]:
    """Return sorted list of supported RAG upload extensions."""

    supported: List[str] = []
    for doc_type, config in RAG_DOCUMENT_TYPE_CONFIG.items():
        if not include_archives and doc_type == "archive":
            continue
        supported.extend(sorted(config.get("extensions", [])))
    return sorted(set(supported))


def _resolve_rag_document_type(suffix: str) -> str:
    """Map a file suffix to configured RAG document type."""

    suffix = (suffix or "").lower().strip(".")
    for doc_type, config in RAG_DOCUMENT_TYPE_CONFIG.items():
        if suffix in config.get("extensions", set()):
            return doc_type
    return "unknown"


def _records_match(left: Iterable[Dict[str, Any]], right: Iterable[Dict[str, Any]]) -> bool:
    """Return ``True`` when two iterable collections of mapping-like rows match."""

    try:
        left_serialised = json.dumps(list(left), sort_keys=True)
        right_serialised = json.dumps(list(right), sort_keys=True)
    except TypeError:
        return False
    return left_serialised == right_serialised


def _build_revenue_stack_chart(summary_df: pd.DataFrame) -> Optional[alt.Chart]:
    """Return a stacked area chart showing annual revenue by category."""

    if summary_df.empty:
        return None

    chart_data = summary_df.copy()
    if "Year" not in chart_data.columns or "Revenue" not in chart_data.columns:
        return None

    chart_data["Year"] = _to_numeric(chart_data["Year"])
    chart_data["Revenue"] = _to_numeric(chart_data["Revenue"])
    chart_data = chart_data.dropna(subset=["Year", "Revenue", "Category"])
    if chart_data.empty:
        return None

    chart_data["Year"] = chart_data["Year"].astype(int)
    unique_years = sorted(chart_data["Year"].unique())

    if len(unique_years) <= 1:
        return (
            alt.Chart(chart_data)
            .mark_bar(opacity=0.85)
            .encode(
                x=alt.X("Category:N", title="Category"),
                y=alt.Y("sum(Revenue):Q", title="Revenue (USD)"),
                color=alt.Color("Category:N", title="Category"),
                tooltip=[
                    alt.Tooltip("Category:N"),
                    alt.Tooltip("Year:O"),
                    alt.Tooltip("Revenue:Q", title="Revenue", format=",.0f"),
                ],
            )
        )

    return (
        alt.Chart(chart_data)
        .mark_area(opacity=0.75)
        .encode(
            x=alt.X("Year:O", title="Year"),
            y=alt.Y("sum(Revenue):Q", title="Revenue (USD)", stack="zero"),
            color=alt.Color("Category:N", title="Category"),
            tooltip=[
                alt.Tooltip("Category:N"),
                alt.Tooltip("Year:O"),
                alt.Tooltip("Revenue:Q", title="Revenue", format=",.0f"),
            ],
        )
    )


def _prepare_cashflow_bridge_frames(cashflow_df: pd.DataFrame) -> Dict[int, pd.DataFrame]:
    """Build per-year cash-flow bridge frames for waterfall visualisations."""

    bridge_frames: Dict[int, pd.DataFrame] = {}
    if cashflow_df.empty or "year" not in cashflow_df.columns:
        return bridge_frames

    numeric_df = cashflow_df.apply(_to_numeric)

    for _, row in numeric_df.iterrows():
        year_value = row.get("year")
        if pd.isna(year_value):
            continue
        year = int(year_value)

        components: List[Tuple[str, float]] = []
        ocf = row.get("operating_cash_flow")
        maint = row.get("maintenance_capex")
        debt_service = row.get("debt_service")
        free_cash_flow = row.get("free_cash_flow")
        present_value = row.get("present_value")

        if not pd.isna(ocf):
            components.append(("Operating cash flow", float(ocf)))
        if not pd.isna(maint):
            components.append(("Maintenance capex", -abs(float(maint))))
        if not pd.isna(debt_service):
            components.append(("Debt service", -abs(float(debt_service))))

        records: List[Dict[str, Any]] = []
        running_total = 0.0

        for label, amount in components:
            start_value = running_total
            running_total += amount
            records.append(
                {
                    "Year": year,
                    "Component": label,
                    "Start": start_value,
                    "End": running_total,
                    "Amount": amount,
                    "Type": "Increase" if amount >= 0 else "Decrease",
                }
            )

        if not pd.isna(free_cash_flow):
            fcf_val = float(free_cash_flow)
            records.append(
                {
                    "Year": year,
                    "Component": "Free cash flow",
                    "Start": 0.0,
                    "End": fcf_val,
                    "Amount": fcf_val,
                    "Type": "Result",
                }
            )

        if not pd.isna(present_value):
            pv_val = float(present_value)
            records.append(
                {
                    "Year": year,
                    "Component": "Present value",
                    "Start": 0.0,
                    "End": pv_val,
                    "Amount": pv_val,
                    "Type": "Result",
                }
            )

        if not records:
            continue

        frame = pd.DataFrame(records)
        frame["Base"] = frame[["Start", "End"]].min(axis=1)
        frame["Cap"] = frame[["Start", "End"]].max(axis=1)
        frame["Mid"] = frame[["Base", "Cap"]].mean(axis=1)
        bridge_frames[year] = frame

    return bridge_frames


def _cashflow_waterfall_chart(frame: pd.DataFrame) -> Optional[alt.Chart]:
    """Return a waterfall chart for a single year's cash-flow bridge."""

    if frame.empty:
        return None

    component_order = [
        label
        for label in [
            "Operating cash flow",
            "Maintenance capex",
            "Debt service",
            "Free cash flow",
            "Present value",
        ]
        if label in frame["Component"].values
    ]

    color_scale = alt.Scale(
        domain=["Increase", "Decrease", "Result"],
        range=["#2ca02c", "#d62728", "#1f77b4"],
    )

    bars = (
        alt.Chart(frame)
        .mark_bar()
        .encode(
            x=alt.X("Component:N", sort=component_order, title="Component"),
            y=alt.Y("Base:Q", title="Cash impact (USD)"),
            y2="Cap:Q",
            color=alt.Color("Type:N", scale=color_scale, title="Movement"),
            tooltip=[
                alt.Tooltip("Component:N"),
                alt.Tooltip("Amount:Q", format=",.0f"),
                alt.Tooltip("Start:Q", title="Start", format=",.0f"),
                alt.Tooltip("End:Q", title="End", format=",.0f"),
            ],
        )
    )

    labels = (
        alt.Chart(frame)
        .mark_text(color="black", dy=-6)
        .encode(
            x=alt.X("Component:N", sort=component_order),
            y=alt.Y("Cap:Q"),
            text=alt.Text("Amount:Q", format=",.0f"),
        )
    )

    return bars + labels


def _combined_leverage_chart(
    dscr_df: pd.DataFrame, coverage_df: pd.DataFrame, leverage_df: pd.DataFrame
) -> Optional[alt.Chart]:
    """Create a combined line chart for DSCR, interest coverage, and leverage metrics."""

    records: List[Dict[str, Any]] = []

    if not dscr_df.empty and {"year", "dscr"}.issubset(dscr_df.columns):
        temp = dscr_df.copy()
        temp["year"] = _to_numeric(temp["year"])
        temp["dscr"] = _to_numeric(temp["dscr"])
        for _, row in temp.dropna(subset=["year", "dscr"]).iterrows():
            records.append(
                {
                    "Year": int(row["year"]),
                    "Metric": "DSCR",
                    "Value": float(row["dscr"]),
                }
            )

    if not coverage_df.empty and "year" in coverage_df.columns:
        temp = coverage_df.copy()
        temp["year"] = _to_numeric(temp["year"])
        for column, label in [
            ("interest_coverage", "Interest coverage"),
        ]:
            if column in temp.columns:
                temp[column] = _to_numeric(temp[column])
                for _, row in temp.dropna(subset=["year", column]).iterrows():
                    records.append(
                        {
                            "Year": int(row["year"]),
                            "Metric": label,
                            "Value": float(row[column]),
                        }
                    )

    if not leverage_df.empty and "year" in leverage_df.columns:
        temp = leverage_df.copy()
        temp["year"] = _to_numeric(temp["year"])
        leverage_column = None
        leverage_label = ""
        if "debt_to_equity" in temp.columns:
            leverage_column = "debt_to_equity"
            leverage_label = "Debt-to-equity"
        elif "debt_ratio" in temp.columns:
            leverage_column = "debt_ratio"
            leverage_label = "Debt ratio"
        if leverage_column:
            temp[leverage_column] = _to_numeric(temp[leverage_column])
            for _, row in temp.dropna(subset=["year", leverage_column]).iterrows():
                records.append(
                    {
                        "Year": int(row["year"]),
                        "Metric": leverage_label,
                        "Value": float(row[leverage_column]),
                    }
                )

    if not records:
        return None

    data = pd.DataFrame(records)
    data = data.dropna(subset=["Year", "Value"])
    if data.empty:
        return None

    data["Year"] = data["Year"].astype(int)

    return (
        alt.Chart(data)
        .mark_line(point=True)
        .encode(
            x=alt.X("Year:O", title="Year"),
            y=alt.Y("Value:Q", title="Ratio"),
            color=alt.Color("Metric:N", title="Metric"),
            tooltip=[
                alt.Tooltip("Metric:N"),
                alt.Tooltip("Year:O"),
                alt.Tooltip("Value:Q", format=",.2f"),
            ],
        )
    )


def _format_currency(value: Any) -> str:
    """Return a compact currency representation for numeric values."""

    if value is None or (isinstance(value, float) and math.isnan(value)):
        return "N/A"
    try:
        numeric = float(value)
    except (TypeError, ValueError):
        return "N/A"
    return f"${numeric:,.0f}"


def _format_ratio(value: Any, digits: int = 2) -> str:
    """Return a fixed precision ratio string for display."""

    if value is None or (isinstance(value, float) and math.isnan(value)):
        return "N/A"
    try:
        numeric = float(value)
    except (TypeError, ValueError):
        return "N/A"
    return f"{numeric:.{digits}f}x"


def _build_income_composition_chart(income_df: pd.DataFrame) -> Optional[alt.Chart]:
    """Render a grouped bar chart for revenue, EBITDA, and net income by year."""

    if income_df.empty or "year" not in income_df.columns:
        return None

    records: List[Dict[str, Any]] = []
    for _, row in income_df.iterrows():
        year = row.get("year")
        for column, label in (
            ("revenue", "Revenue"),
            ("ebitda", "EBITDA"),
            ("net_income", "Net income"),
        ):
            value = row.get(column)
            if year is None or pd.isna(year) or value is None or pd.isna(value):
                continue
            records.append(
                {
                    "Year": int(float(year)),
                    "Metric": label,
                    "Value": float(value),
                }
            )

    if not records:
        return None

    chart_df = pd.DataFrame(records)
    return (
        alt.Chart(chart_df)
        .mark_bar()
        .encode(
            x=alt.X("Year:O", title="Year"),
            y=alt.Y("Value:Q", title="USD"),
            color=alt.Color("Metric:N", title="Metric"),
            xOffset="Metric:N",
            tooltip=[
                alt.Tooltip("Year:O"),
                alt.Tooltip("Metric:N"),
                alt.Tooltip("Value:Q", format=",.0f"),
            ],
        )
    )


def _build_business_plan_markdown(
    scenario: str,
    assumptions: Assumptions,
    valuation: Dict[str, Any],
    annual_df: pd.DataFrame,
    dscr_df: pd.DataFrame,
    break_even_df: pd.DataFrame,
    monte_carlo_summary_df: pd.DataFrame,
) -> str:
    """Build comprehensive business plan narrative text from current model outputs."""

    annual_row = annual_df.iloc[0].to_dict() if not annual_df.empty else {}
    start_year = int(getattr(assumptions, "production_start_year", 0) or 0)
    horizon = int(getattr(assumptions, "production_horizon_years", 0) or 0)
    end_year = start_year + max(horizon - 1, 0)

    avg_dscr = dscr_df["dscr"].mean() if not dscr_df.empty and "dscr" in dscr_df.columns else None
    avg_break_even_price = (
        break_even_df["break_even_price_per_kg"].mean()
        if not break_even_df.empty and "break_even_price_per_kg" in break_even_df.columns
        else None
    )

    npv = valuation.get("npv")
    irr = valuation.get("irr")
    irr_text = f"{float(irr):.2%}" if isinstance(irr, (int, float)) else "N/A"
    payback = valuation.get("payback_period_years")
    annual_revenue = annual_row.get("revenue")
    annual_ebitda = annual_row.get("ebitda")
    annual_net_income = annual_row.get("net_income")

    monte_carlo_p5 = None
    monte_carlo_p50 = None
    monte_carlo_p95 = None
    if not monte_carlo_summary_df.empty:
        row = monte_carlo_summary_df.iloc[0]
        monte_carlo_p5 = row.get("npv_p5")
        monte_carlo_p50 = row.get("npv_p50")
        monte_carlo_p95 = row.get("npv_p95")

    return f"""
### 1) Executive Summary
The **{scenario}** strategy for the broiler project is designed for a planning horizon from **{start_year} to {end_year}**.  
The model estimates a project NPV of **{_format_currency(npv)}** and an IRR of **{irr_text}**, indicating expected value creation under current assumptions.

### 2) Production & Operating Plan
- Cycles per year: **{assumptions.cycles_per_year}**
- Birds per cycle: **{assumptions.birds_per_cycle:,}**
- Mortality rate: **{assumptions.mortality_rate:.2%}**
- Feed conversion ratio: **{assumptions.feed_conversion_ratio:.2f}**

Operational execution should prioritize biosecurity discipline, feed-efficiency optimization, and labor/energy productivity to stabilize margin delivery.

### 3) Financial Performance Plan
- Annual revenue (current run): **{_format_currency(annual_revenue)}**
- Annual EBITDA (current run): **{_format_currency(annual_ebitda)}**
- Annual net income (current run): **{_format_currency(annual_net_income)}**
- Average DSCR: **{_format_ratio(avg_dscr)}**
- Payback period: **{payback if payback is not None else 'N/A'} years**

The finance plan should maintain covenant headroom with DSCR buffers and a liquidity reserve calibrated to debt-service and working-capital variability.

### 4) Market, Pricing, and Break-even Strategy
- Base live price per kg: **{_format_currency(assumptions.live_price_per_kg)}**
- Price growth assumption: **{assumptions.price_growth:.2%}**
- Average break-even price per kg: **{_format_currency(avg_break_even_price)}**

Commercial strategy should blend contracted offtake and spot exposure to protect downside while preserving upside during favorable price cycles.

### 5) Risk Assessment & Mitigation
- Monte Carlo NPV P5: **{_format_currency(monte_carlo_p5)}**
- Monte Carlo NPV P50: **{_format_currency(monte_carlo_p50)}**
- Monte Carlo NPV P95: **{_format_currency(monte_carlo_p95)}**

Risk controls should include feed-procurement hedging, contingency mortality protocols, and refinancing reviews to reduce free-cash-flow volatility.

### 6) Implementation Roadmap
1. **Year {start_year}:** commissioning, supplier onboarding, and operating rhythm stabilization.  
2. **Years {start_year + 1}-{end_year}:** throughput optimization, cost benchmarking, and selective capacity/process upgrades.  
3. **Annual cadence:** quarterly KPI reviews (margin, DSCR, mortality, FCR), scenario stress tests, and pricing strategy updates.
""".strip()


def _generate_investor_recommendations(
    assumptions: Assumptions,
    valuation: Dict[str, Any],
    annual_df: pd.DataFrame,
    dscr_df: pd.DataFrame,
    break_even_df: pd.DataFrame,
) -> List[str]:
    """Generate practical recommendations that improve investor appeal."""

    recommendations: List[str] = []
    npv = valuation.get("npv")
    irr = valuation.get("irr")
    payback = valuation.get("payback_period_years")

    avg_dscr = (
        float(dscr_df["dscr"].mean())
        if not dscr_df.empty and "dscr" in dscr_df.columns
        else None
    )
    avg_break_even_price = (
        float(break_even_df["break_even_price_per_kg"].mean())
        if not break_even_df.empty and "break_even_price_per_kg" in break_even_df.columns
        else None
    )

    annual_revenue = None
    if not annual_df.empty and "revenue" in annual_df.columns:
        revenue = annual_df.iloc[0].get("revenue")
        annual_revenue = float(revenue) if pd.notna(revenue) else None

    if isinstance(npv, (int, float)) and npv > 0:
        recommendations.append(
            f"Lead with the positive NPV case (**{_format_currency(npv)}**) and show a "
            "base/downside/upside sensitivity table in investor materials."
        )
    else:
        recommendations.append(
            "Rework the investment thesis by improving price/cost assumptions and "
            "showing a credible path to positive NPV before investor outreach."
        )

    if isinstance(irr, (int, float)):
        if irr >= 0.18:
            recommendations.append(
                f"Position IRR (**{irr:.2%}**) as a strong return profile versus agri-sector "
                "benchmarks while documenting key assumptions transparently."
            )
        else:
            recommendations.append(
                f"IRR (**{irr:.2%}**) is modest; improve attractiveness by combining operating "
                "efficiency initiatives with financing optimization."
            )

    if avg_dscr is not None:
        if avg_dscr >= 1.5:
            recommendations.append(
                f"Highlight debt-service resilience (average DSCR **{avg_dscr:.2f}x**) and "
                "include covenant headroom charts to reassure lenders and equity partners."
            )
        else:
            recommendations.append(
                f"Strengthen bankability by improving DSCR (currently **{avg_dscr:.2f}x**) "
                "through lower leverage or staged capex."
            )

    if avg_break_even_price is not None:
        price_gap = assumptions.live_price_per_kg - avg_break_even_price
        if price_gap > 0:
            recommendations.append(
                f"Show pricing buffer: current live price exceeds break-even by "
                f"**{_format_currency(price_gap)} per kg**, reinforcing downside protection."
            )
        else:
            recommendations.append(
                "Current live price is at/below break-even; prioritize cost-out actions and "
                "offtake contract repricing before a fundraise."
            )

    if isinstance(payback, (int, float)):
        recommendations.append(
            f"Present a milestone-based capital recovery timeline with a modeled payback of "
            f"**{payback:.1f} years** and quarterly KPI checkpoints."
        )

    if annual_revenue is not None:
        recommendations.append(
            f"Package the model with a one-page investment memo anchored on expected annual "
            f"revenue (**{_format_currency(annual_revenue)}**), unit economics, and risk controls."
        )

    recommendations.append(
        "Increase investor trust by publishing a model governance pack: assumptions log, "
        "version history, and scenario-testing methodology."
    )

    return recommendations


def _confidence_band(label: str, value: float) -> str:
    """Return confidence level for selected assumptions."""

    if label == "Mortality rate":
        return "High" if value <= 0.05 else "Medium" if value <= 0.08 else "Low"
    if label == "Feed conversion ratio":
        return "High" if value <= 1.65 else "Medium" if value <= 1.85 else "Low"
    if label == "Feed cost per kg":
        return "High" if value <= 0.50 else "Medium" if value <= 0.65 else "Low"
    if label == "Live price per kg":
        return "High" if value >= 1.8 else "Medium" if value >= 1.5 else "Low"
    return "Medium"


def _build_assumption_confidence_frame(assumptions: Assumptions) -> pd.DataFrame:
    """Create confidence ratings for key investment assumptions."""

    items = [
        ("Live price per kg", float(assumptions.live_price_per_kg)),
        ("Feed conversion ratio", float(assumptions.feed_conversion_ratio)),
        ("Mortality rate", float(assumptions.mortality_rate)),
        ("Feed cost per kg", float(assumptions.feed_cost_per_kg)),
    ]
    rows: List[Dict[str, Any]] = []
    for label, value in items:
        rows.append(
            {
                "Assumption": label,
                "Value": value,
                "Confidence": _confidence_band(label, value),
            }
        )
    return pd.DataFrame(rows)


def _build_investor_scorecard(
    valuation: Dict[str, Any],
    dscr_df: pd.DataFrame,
    break_even_df: pd.DataFrame,
    assumptions: Assumptions,
    benchmarks: Dict[str, float],
) -> pd.DataFrame:
    """Build traffic-light scorecard for investors."""

    irr = valuation.get("irr")
    payback = valuation.get("payback_period_years")
    npv = valuation.get("npv")
    avg_dscr = (
        float(dscr_df["dscr"].mean())
        if not dscr_df.empty and "dscr" in dscr_df.columns
        else None
    )
    break_even_price = (
        float(break_even_df["break_even_price_per_kg"].mean())
        if not break_even_df.empty and "break_even_price_per_kg" in break_even_df.columns
        else None
    )
    price_buffer = (
        float(assumptions.live_price_per_kg) - break_even_price
        if break_even_price is not None
        else None
    )

    irr_threshold = float(benchmarks.get("target_irr", 0.18))
    dscr_threshold = float(benchmarks.get("min_dscr", 1.5))
    payback_threshold = float(benchmarks.get("max_payback_years", 6.0))

    metrics = [
        ("NPV", npv, "positive"),
        ("IRR", irr, "high"),
        ("Average DSCR", avg_dscr, "high"),
        ("Payback (years)", payback, "low"),
        ("Price buffer vs break-even", price_buffer, "positive"),
    ]
    rows: List[Dict[str, Any]] = []
    for metric, value, direction in metrics:
        status = "🟡 Watch"
        if value is None or (isinstance(value, float) and math.isnan(value)):
            status = "⚪ N/A"
        elif direction == "positive":
            status = "🟢 Strong" if float(value) > 0 else "🔴 Weak"
        elif direction == "high":
            threshold = irr_threshold if metric == "IRR" else dscr_threshold
            status = "🟢 Strong" if float(value) >= threshold else "🔴 Weak"
        elif direction == "low":
            status = "🟢 Strong" if float(value) <= payback_threshold else "🔴 Weak"

        rows.append({"Metric": metric, "Value": value, "Status": status})

    return pd.DataFrame(rows)


def _build_ic_pack_markdown(
    scenario: str,
    scorecard_df: pd.DataFrame,
    plan_markdown: str,
    recommendations: List[str],
) -> str:
    """Create an investment-committee text pack for download."""

    score_lines = []
    for _, row in scorecard_df.iterrows():
        score_lines.append(
            f"- {row.get('Metric')}: {row.get('Status')} ({row.get('Value')})"
        )
    recommendation_lines = [f"- {item}" for item in recommendations]
    return "\n".join(
        [
            f"# Investment Committee Pack - {scenario}",
            "",
            "## Investor Scorecard",
            *score_lines,
            "",
            "## Comprehensive Plan",
            plan_markdown,
            "",
            "## Recommended Actions",
            *recommendation_lines,
        ]
    )


def _generate_business_plan_excel_bytes(
    scenario: str,
    plan_markdown: str,
    recommendations: List[str],
    scorecard_df: pd.DataFrame,
    confidence_df: pd.DataFrame,
    citations: Optional[List[Dict[str, str]]] = None,
    financial_tables: Optional[Dict[str, pd.DataFrame]] = None,
    chart_payloads: Optional[Dict[str, Dict[str, Any]]] = None,
) -> bytes:
    """Create a business-plan Excel workbook export."""

    output = BytesIO()
    try:
        import xlsxwriter  # type: ignore  # noqa: F401

        engine = "xlsxwriter"
    except ImportError:
        try:
            import openpyxl  # type: ignore  # noqa: F401

            engine = "openpyxl"
        except ImportError as exc:
            raise RuntimeError("Missing Excel writer dependency") from exc

    with pd.ExcelWriter(output, engine=engine) as writer:
        pd.DataFrame([{"Section": "Business Plan", "Content": plan_markdown}]).to_excel(
            writer, sheet_name="Business Plan", index=False
        )
        pd.DataFrame({"Recommendation": recommendations}).to_excel(
            writer, sheet_name="Recommendations", index=False
        )
        scorecard_df.to_excel(writer, sheet_name="Scorecard", index=False)
        confidence_df.to_excel(writer, sheet_name="Confidence Bands", index=False)
        if citations:
            pd.DataFrame(citations).to_excel(
                writer, sheet_name="RAG Citations", index=False
            )
        if financial_tables:
            for sheet_name, table in financial_tables.items():
                safe_sheet_name = re.sub(r"[^A-Za-z0-9 ]+", "", sheet_name).strip()
                safe_sheet_name = safe_sheet_name[:31] or "Model Data"
                export_df = table.copy()
                for column in export_df.columns:
                    if export_df[column].dtype == object:
                        export_df[column] = export_df[column].apply(
                            lambda value: json.dumps(value)
                            if isinstance(value, (dict, list, tuple, set))
                            else value
                        )
                export_df.to_excel(writer, sheet_name=safe_sheet_name, index=False)
        if chart_payloads:
            chart_index_rows: List[Dict[str, str]] = []
            for chart_name, payload in chart_payloads.items():
                safe_name = re.sub(r"[^A-Za-z0-9 ]+", "", chart_name).strip() or "Chart"
                data_sheet = f"{safe_name} Data"[:31]
                spec_sheet = f"{safe_name} Spec"[:31]
                data_frame = payload.get("data")
                if isinstance(data_frame, pd.DataFrame):
                    data_frame.to_excel(writer, sheet_name=data_sheet, index=False)
                spec = payload.get("spec")
                if spec:
                    pd.DataFrame(
                        [{"chart": chart_name, "vega_lite_spec_json": json.dumps(spec)}]
                    ).to_excel(writer, sheet_name=spec_sheet, index=False)
                chart_index_rows.append(
                    {"Chart": chart_name, "Data Sheet": data_sheet, "Spec Sheet": spec_sheet}
                )
            if chart_index_rows:
                pd.DataFrame(chart_index_rows).to_excel(
                    writer, sheet_name="Charts Index", index=False
                )

        if engine == "xlsxwriter":
            writer.book.set_properties(
                {"title": f"Business Plan - {scenario}", "subject": "Investor pack"}
            )

    output.seek(0)
    return output.getvalue()


def _generate_business_plan_csv_zip_bytes(
    scenario: str,
    plan_markdown: str,
    recommendations: List[str],
    scorecard_df: pd.DataFrame,
    confidence_df: pd.DataFrame,
    citations: Optional[List[Dict[str, str]]] = None,
    financial_tables: Optional[Dict[str, pd.DataFrame]] = None,
    chart_payloads: Optional[Dict[str, Dict[str, Any]]] = None,
) -> bytes:
    """Create a ZIP of CSV/JSON artifacts when Excel engines are unavailable."""

    buffer = BytesIO()
    with zipfile.ZipFile(buffer, mode="w", compression=zipfile.ZIP_DEFLATED) as archive:
        archive.writestr(
            "README.txt",
            (
                f"Business plan export for scenario: {scenario}\n"
                "Excel engines were unavailable, so this ZIP includes CSV and JSON artifacts.\n"
            ),
        )
        archive.writestr("business_plan.md", plan_markdown)
        archive.writestr(
            "recommendations.csv",
            pd.DataFrame({"Recommendation": recommendations}).to_csv(index=False),
        )
        archive.writestr("scorecard.csv", scorecard_df.to_csv(index=False))
        archive.writestr("confidence_bands.csv", confidence_df.to_csv(index=False))
        if citations:
            archive.writestr("rag_citations.csv", pd.DataFrame(citations).to_csv(index=False))
        if financial_tables:
            for name, table in financial_tables.items():
                safe_name = re.sub(r"[^a-z0-9]+", "_", name.lower()).strip("_")
                archive.writestr(
                    f"tables/{safe_name or 'model_data'}.csv", table.to_csv(index=False)
                )
        if chart_payloads:
            index_rows: List[Dict[str, str]] = []
            for chart_name, payload in chart_payloads.items():
                safe_name = re.sub(r"[^a-z0-9]+", "_", chart_name.lower()).strip("_")
                data_frame = payload.get("data")
                if isinstance(data_frame, pd.DataFrame):
                    archive.writestr(
                        f"charts/{safe_name or 'chart'}_data.csv",
                        data_frame.to_csv(index=False),
                    )
                spec = payload.get("spec")
                if spec:
                    archive.writestr(
                        f"charts/{safe_name or 'chart'}_spec.json",
                        json.dumps(spec, indent=2),
                    )
                index_rows.append(
                    {
                        "chart": chart_name,
                        "data_file": f"charts/{safe_name or 'chart'}_data.csv",
                        "spec_file": f"charts/{safe_name or 'chart'}_spec.json",
                    }
                )
            if index_rows:
                archive.writestr(
                    "charts/index.csv", pd.DataFrame(index_rows).to_csv(index=False)
                )
    buffer.seek(0)
    return buffer.getvalue()


def _generate_business_plan_docx_bytes(
    title: str,
    plan_markdown: str,
    recommendations: List[str],
    citations: Optional[List[Dict[str, str]]] = None,
) -> bytes:
    """Create a Word document export for the generated business plan."""

    try:
        from docx import Document  # type: ignore
    except ImportError as exc:
        raise RuntimeError("Word export requires `python-docx`.") from exc

    document = Document()
    document.add_heading(title, level=1)
    for paragraph in plan_markdown.split("\n\n"):
        stripped = paragraph.strip()
        if stripped:
            document.add_paragraph(stripped)

    document.add_heading("Investor Recommendations", level=2)
    for item in recommendations:
        document.add_paragraph(item, style="List Bullet")
    if citations:
        document.add_heading("RAG Source Citations", level=2)
        for citation in citations:
            document.add_paragraph(
                f"{citation.get('section')}: {citation.get('source')} ({citation.get('chunk_id')})",
                style="List Bullet",
            )

    output = BytesIO()
    document.save(output)
    output.seek(0)
    return output.getvalue()


def _generate_business_plan_pdf_bytes(
    title: str,
    plan_markdown: str,
    recommendations: List[str],
    citations: Optional[List[Dict[str, str]]] = None,
) -> bytes:
    """Create a PDF export for the generated business plan."""

    try:
        from reportlab.lib.pagesizes import letter  # type: ignore
        from reportlab.pdfgen import canvas  # type: ignore
    except ImportError as exc:
        raise RuntimeError("PDF export requires `reportlab`.") from exc

    output = BytesIO()
    pdf = canvas.Canvas(output, pagesize=letter)
    width, height = letter
    y = height - 40

    def _write_line(text: str, *, bold: bool = False) -> None:
        nonlocal y
        if y < 50:
            pdf.showPage()
            y = height - 40
        pdf.setFont("Helvetica-Bold" if bold else "Helvetica", 10)
        pdf.drawString(40, y, text[:120])
        y -= 14

    _write_line(title, bold=True)
    _write_line("")
    for line in plan_markdown.splitlines():
        _write_line(line)
    _write_line("")
    _write_line("Investor Recommendations", bold=True)
    for item in recommendations:
        _write_line(f"- {item}")
    if citations:
        _write_line("")
        _write_line("RAG Source Citations", bold=True)
        for citation in citations:
            _write_line(
                f"- {citation.get('section')}: {citation.get('source')} ({citation.get('chunk_id')})"
            )

    pdf.save()
    output.seek(0)
    return output.getvalue()


def _extract_text_from_upload(uploaded_file: Any) -> Tuple[str, str, Optional[str]]:
    """Extract text from uploaded files and return (text, parser_used, error_message)."""

    def _extract_docx_via_zip(data: bytes) -> str:
        with zipfile.ZipFile(BytesIO(data)) as archive:
            xml_payload = archive.read("word/document.xml")
        root = ET.fromstring(xml_payload)
        texts = [node.text for node in root.iter() if node.tag.endswith("}t") and node.text]
        return "\n".join(texts).strip()

    def _extract_xlsx_via_zip(data: bytes) -> str:
        with zipfile.ZipFile(BytesIO(data)) as archive:
            shared_strings: List[str] = []
            if "xl/sharedStrings.xml" in archive.namelist():
                shared_root = ET.fromstring(archive.read("xl/sharedStrings.xml"))
                shared_strings = [
                    "".join(t.text or "" for t in si.iter() if t.tag.endswith("}t"))
                    for si in shared_root.iter()
                    if si.tag.endswith("}si")
                ]
            sheet_files = [
                name for name in archive.namelist() if name.startswith("xl/worksheets/sheet")
            ]
            lines: List[str] = []
            for sheet_name in sorted(sheet_files):
                root = ET.fromstring(archive.read(sheet_name))
                lines.append(f"[Sheet: {sheet_name.split('/')[-1]}]")
                for row in root.iter():
                    if not row.tag.endswith("}row"):
                        continue
                    values: List[str] = []
                    for cell in row:
                        if not cell.tag.endswith("}c"):
                            continue
                        cell_type = cell.attrib.get("t")
                        value_node = next((c for c in cell if c.tag.endswith("}v")), None)
                        if value_node is None or value_node.text is None:
                            continue
                        value_text = value_node.text
                        if cell_type == "s":
                            try:
                                idx = int(value_text)
                                value_text = shared_strings[idx] if idx < len(shared_strings) else value_text
                            except Exception:
                                pass
                        values.append(value_text)
                    if values:
                        lines.append(" | ".join(values))
            return "\n".join(lines).strip()

    def _extract_pptx_via_zip(data: bytes) -> str:
        with zipfile.ZipFile(BytesIO(data)) as archive:
            slide_files = [
                name for name in archive.namelist() if name.startswith("ppt/slides/slide")
            ]
            lines: List[str] = []
            for slide_file in sorted(slide_files):
                lines.append(f"[{slide_file.split('/')[-1]}]")
                root = ET.fromstring(archive.read(slide_file))
                texts = [
                    node.text
                    for node in root.iter()
                    if node.tag.endswith("}t") and node.text and node.text.strip()
                ]
                lines.extend(texts)
            return "\n".join(lines).strip()

    file_name = getattr(uploaded_file, "name", "document")
    suffix = file_name.lower().split(".")[-1] if "." in file_name else ""
    file_bytes = uploaded_file.getvalue()
    doc_type = _resolve_rag_document_type(suffix)

    if doc_type == "text":
        return _decode_bytes_to_text(file_bytes), "plain-text", None

    if doc_type == "delimited":
        delimiter = "\t" if suffix == "tsv" else ","
        try:
            decoded = _decode_bytes_to_text(file_bytes)
            lines = decoded.splitlines()
            if not lines:
                return "", "csv", "File is empty"
            return decoded, f"delimited({delimiter})", None
        except Exception as exc:
            return "", "delimited", str(exc)

    if doc_type == "json":
        try:
            obj = json.loads(_decode_bytes_to_text(file_bytes))
            return json.dumps(obj, indent=2), "json", None
        except Exception as exc:
            return "", "json", str(exc)

    if doc_type == "markup":
        text = _decode_bytes_to_text(file_bytes)
        cleaned = re.sub(r"<[^>]+>", " ", text)
        cleaned = re.sub(r"\s+", " ", cleaned).strip()
        return cleaned, "markup-strip", None

    if doc_type == "pdf":
        pdf_failures: List[str] = []

        def _record_pdf_failure(step: str, exc: Optional[Exception] = None) -> None:
            if exc is None:
                pdf_failures.append(step)
                return
            message = str(exc).strip()
            reason = f"{type(exc).__name__}: {message}" if message else type(exc).__name__
            pdf_failures.append(f"{step} ({reason})")

        reader_cls = None
        reader_label = "pypdf/PyPDF2"
        try:
            from pypdf import PdfReader as _PdfReader  # type: ignore

            reader_cls = _PdfReader
            reader_label = "pypdf"
        except ImportError:
            try:
                from PyPDF2 import PdfReader as _PdfReader  # type: ignore

                reader_cls = _PdfReader
                reader_label = "PyPDF2"
            except ImportError:
                _record_pdf_failure("pypdf/PyPDF2 unavailable")
                reader_cls = None
        if reader_cls is not None:
            try:
                reader = reader_cls(BytesIO(file_bytes))
                if getattr(reader, "is_encrypted", False):
                    try:
                        reader.decrypt("")  # type: ignore[attr-defined]
                    except Exception:
                        pass
                pages = [page.extract_text() or "" for page in reader.pages]
                extracted = "\n".join(pages).strip()
                if extracted:
                    return extracted, getattr(reader_cls, "__module__", "pdf-reader"), None
                _record_pdf_failure(f"{reader_label} returned empty text")
            except Exception as exc:
                _record_pdf_failure(f"{reader_label} failed", exc)

        try:
            import pdfplumber  # type: ignore

            with pdfplumber.open(BytesIO(file_bytes)) as pdf:
                pages = [(page.extract_text() or "") for page in pdf.pages]
            extracted = "\n".join(pages).strip()
            if extracted:
                return extracted, "pdfplumber", None
            _record_pdf_failure("pdfplumber returned empty text")
        except ImportError:
            _record_pdf_failure("pdfplumber unavailable")
        except Exception as exc:
            _record_pdf_failure("pdfplumber failed", exc)

        try:
            import fitz  # type: ignore

            document = fitz.open(stream=file_bytes, filetype="pdf")
            pages = [page.get_text("text") or "" for page in document]
            extracted = "\n".join(pages).strip()
            if extracted:
                return extracted, "pymupdf", None
            _record_pdf_failure("PyMuPDF returned empty text")
        except ImportError:
            _record_pdf_failure("PyMuPDF unavailable")
        except Exception as exc:
            _record_pdf_failure("PyMuPDF failed", exc)

        try:
            from pdfminer.high_level import extract_text as _pdfminer_extract_text  # type: ignore

            extracted = (_pdfminer_extract_text(BytesIO(file_bytes)) or "").strip()
            if extracted:
                return extracted, "pdfminer.six", None
            _record_pdf_failure("pdfminer.six returned empty text")
        except ImportError:
            _record_pdf_failure("pdfminer.six unavailable")
        except Exception as exc:
            _record_pdf_failure("pdfminer.six failed", exc)

        if shutil.which("pdftotext"):
            try:
                with tempfile.NamedTemporaryFile(suffix=".pdf") as src_file, tempfile.NamedTemporaryFile(
                    suffix=".txt"
                ) as out_file:
                    src_file.write(file_bytes)
                    src_file.flush()
                    subprocess.run(
                        ["pdftotext", "-layout", src_file.name, out_file.name],
                        check=True,
                        capture_output=True,
                        text=True,
                    )
                    out_file.seek(0)
                    extracted = _decode_bytes_to_text(out_file.read()).strip()
                    if extracted:
                        return extracted, "pdftotext", None
                    _record_pdf_failure("pdftotext returned empty text")
            except Exception as exc:
                _record_pdf_failure("pdftotext failed", exc)
        else:
            _record_pdf_failure("pdftotext binary unavailable")

        if shutil.which("tesseract"):
            try:
                from pdf2image import convert_from_bytes  # type: ignore
                import pytesseract  # type: ignore

                images = convert_from_bytes(file_bytes, dpi=220)
                ocr_pages = [(pytesseract.image_to_string(img) or "").strip() for img in images]
                extracted = "\n".join([text for text in ocr_pages if text]).strip()
                if extracted:
                    return extracted, "ocr-pytesseract(pdf2image)", None
                _record_pdf_failure("OCR via pdf2image returned empty text")
            except ImportError:
                _record_pdf_failure("OCR via pdf2image unavailable")
            except Exception as exc:
                _record_pdf_failure("OCR via pdf2image failed", exc)

            try:
                import fitz  # type: ignore
                from PIL import Image  # type: ignore
                import pytesseract  # type: ignore

                doc = fitz.open(stream=file_bytes, filetype="pdf")
                ocr_pages: List[str] = []
                for page in doc:
                    pix = page.get_pixmap(dpi=220)
                    image = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    ocr_pages.append((pytesseract.image_to_string(image) or "").strip())
                extracted = "\n".join([text for text in ocr_pages if text]).strip()
                if extracted:
                    return extracted, "ocr-pytesseract(pymupdf)", None
                _record_pdf_failure("OCR via PyMuPDF+Pillow returned empty text")
            except ImportError:
                _record_pdf_failure("OCR via PyMuPDF+Pillow unavailable")
            except Exception as exc:
                _record_pdf_failure("OCR via PyMuPDF+Pillow failed", exc)
        else:
            _record_pdf_failure("tesseract binary unavailable")

        # Dependency-free fallback: recover readable text segments directly from PDF bytes.
        # This is not as accurate as parser libraries, but it prevents hard indexing failure
        # when optional PDF dependencies are unavailable in a deployment environment.
        try:
            decoded = file_bytes.decode("latin-1", errors="ignore")
            decoded = decoded.replace("\x00", " ")
            decoded = re.sub(r"\s+", " ", decoded)
            text_candidates = re.findall(r"[A-Za-z][A-Za-z0-9 ,.;:()\\-_/]{24,}", decoded)
            if text_candidates:
                extracted = "\n".join(text_candidates[:300]).strip()
                if extracted:
                    return extracted, "pdf-byte-fallback", None
            _record_pdf_failure("pdf-byte-fallback returned empty text")
        except Exception as exc:
            _record_pdf_failure("pdf-byte-fallback failed", exc)

        details = "; ".join(pdf_failures[:6]) if pdf_failures else "No parser diagnostics available"
        return "", "pdf", (
            "No PDF parser/OCR could extract text "
            "(pypdf/PyPDF2/pdfplumber/PyMuPDF/pdfminer/pdftotext; OCR requires tesseract + pdf2image or PyMuPDF+Pillow+pytesseract). "
            f"Details: {details}"
        )

    if doc_type == "word":
        try:
            from docx import Document  # type: ignore
        except ImportError:
            try:
                extracted = _extract_docx_via_zip(file_bytes)
                if extracted:
                    return extracted, "docx-zipxml", None
            except Exception:
                pass
            return "", "docx", "python-docx is not installed and ZIP-XML fallback failed"
        try:
            document = Document(BytesIO(file_bytes))
            text = "\n".join(p.text for p in document.paragraphs if p.text.strip())
            if text.strip():
                return text, "python-docx", None
        except Exception:
            pass
        try:
            extracted = _extract_docx_via_zip(file_bytes)
            if extracted:
                return extracted, "docx-zipxml", None
        except Exception as exc:
            return "", "docx", str(exc)
        return "", "docx", "Could not extract text from DOCX"

    if doc_type == "excel" and suffix in {"xlsx", "xlsm", "xltx", "xltm"}:
        try:
            from openpyxl import load_workbook  # type: ignore

            workbook = load_workbook(BytesIO(file_bytes), data_only=True, read_only=True)
            lines: List[str] = []
            for sheet in workbook.worksheets:
                lines.append(f"[Sheet: {sheet.title}]")
                for row in sheet.iter_rows(values_only=True):
                    values = [str(cell) for cell in row if cell is not None and str(cell).strip()]
                    if values:
                        lines.append(" | ".join(values))
            extracted = "\n".join(lines).strip()
            if extracted:
                return extracted, "openpyxl", None
        except Exception as exc:
            openpyxl_error = str(exc)
        else:
            openpyxl_error = ""
        try:
            extracted = _extract_xlsx_via_zip(file_bytes)
            if extracted:
                return extracted, "xlsx-zipxml", None
        except Exception:
            pass
        return "", "xlsx", openpyxl_error or "Could not extract text from XLSX"

    if doc_type == "excel" and suffix == "xls":
        try:
            import pandas as _pd  # type: ignore

            workbook = _pd.read_excel(BytesIO(file_bytes), sheet_name=None)
            lines: List[str] = []
            for sheet_name, frame in workbook.items():
                lines.append(f"[Sheet: {sheet_name}]")
                for _, row in frame.fillna("").astype(str).head(500).iterrows():
                    row_values = [val for val in row.tolist() if str(val).strip()]
                    if row_values:
                        lines.append(" | ".join(row_values))
            extracted = "\n".join(lines).strip()
            if extracted:
                return extracted, "pandas-xls", None
        except Exception as exc:
            return "", "xls", str(exc)
        return "", "xls", "Could not extract text from XLS"

    if doc_type == "pptx":
        try:
            from pptx import Presentation  # type: ignore

            prs = Presentation(BytesIO(file_bytes))
            lines = []
            for idx, slide in enumerate(prs.slides, start=1):
                lines.append(f"[Slide {idx}]")
                for shape in slide.shapes:
                    text = getattr(shape, "text", "")
                    if text and text.strip():
                        lines.append(text.strip())
            extracted = "\n".join(lines).strip()
            if extracted:
                return extracted, "python-pptx", None
        except Exception as exc:
            pptx_error = str(exc)
        else:
            pptx_error = ""
        try:
            extracted = _extract_pptx_via_zip(file_bytes)
            if extracted:
                return extracted, "pptx-zipxml", None
        except Exception:
            pass
        return "", "pptx", pptx_error or "Could not extract text from PPTX"

    fallback = _decode_bytes_to_text(file_bytes).strip()
    if fallback:
        return fallback, "fallback-utf8", None
    return "", "unknown", "Unsupported or binary file format with no text parser available."


def _expand_uploaded_documents(uploaded_docs: Optional[List[Any]]) -> List[Dict[str, Any]]:
    """Expand uploaded files, including ZIP archives, into parseable payloads."""

    expanded: List[Dict[str, Any]] = []
    supported_suffixes = set(_rag_supported_extensions(include_archives=False))
    for uploaded in uploaded_docs or []:
        name = getattr(uploaded, "name", "document")
        file_bytes = uploaded.getvalue()
        suffix = name.lower().split(".")[-1] if "." in name else ""
        if suffix != "zip":
            expanded.append({"name": name, "payload": file_bytes, "supported": suffix in supported_suffixes})
            continue
        try:
            with zipfile.ZipFile(BytesIO(file_bytes)) as archive:
                members = [
                    m
                    for m in archive.namelist()
                    if not m.endswith("/") and not m.startswith("__MACOSX/")
                ]
                for member in members:
                    member_suffix = member.lower().split(".")[-1] if "." in member else ""
                    expanded.append(
                        {
                            "name": f"{name}:{member}",
                            "payload": archive.read(member),
                            "supported": member_suffix in supported_suffixes,
                        }
                    )
        except Exception:
            expanded.append({"name": name, "payload": file_bytes, "supported": False})
    return expanded


def _extract_text_from_payload(file_name: str, payload: bytes) -> Tuple[str, str, Optional[str]]:
    """Parse text from a raw file payload by adapting to upload parser interface."""

    class _UploadedPayload:
        def __init__(self, name: str, data: bytes) -> None:
            self.name = name
            self._data = data

        def getvalue(self) -> bytes:
            return self._data

    return _extract_text_from_upload(_UploadedPayload(file_name, payload))


def _chunk_document_text(
    source_name: str, text: str, chunk_size: int = 1200, overlap: int = 150
) -> List[Dict[str, Any]]:
    """Create overlapping chunks from a source document."""

    clean_text = re.sub(r"\s+", " ", text).strip()
    if not clean_text:
        return []

    chunks: List[Dict[str, Any]] = []
    start = 0
    idx = 1
    while start < len(clean_text):
        end = min(start + chunk_size, len(clean_text))
        excerpt = clean_text[start:end].strip()
        if excerpt:
            chunks.append(
                {
                    "chunk_id": f"{source_name}::chunk_{idx}",
                    "source": source_name,
                    "text": excerpt,
                }
            )
            idx += 1
        if end >= len(clean_text):
            break
        start = max(0, end - overlap)

    return chunks


def _simple_tokenise(text: str) -> List[str]:
    """Tokenise text for lightweight lexical retrieval."""

    return re.findall(r"[a-zA-Z0-9]+", text.lower())


def _retrieve_evidence_for_sections(
    chunks: List[Dict[str, Any]],
    section_queries: Dict[str, str],
    top_k: int = 3,
) -> Dict[str, List[Dict[str, Any]]]:
    """Return top-k chunk matches per section using lexical overlap scoring."""

    if not chunks:
        return {section: [] for section in section_queries}

    chunk_payload: List[Tuple[Dict[str, Any], set[str]]] = []
    for chunk in chunks:
        tokens = set(_simple_tokenise(chunk.get("text", "")))
        chunk_payload.append((chunk, tokens))

    section_matches: Dict[str, List[Dict[str, Any]]] = {}
    for section, query in section_queries.items():
        query_tokens = set(_simple_tokenise(query))
        scored: List[Tuple[float, Dict[str, Any]]] = []
        for chunk, tokens in chunk_payload:
            if not tokens:
                continue
            overlap = len(query_tokens & tokens)
            if overlap <= 0:
                continue
            score = overlap / max(len(query_tokens), 1)
            scored.append((score, chunk))
        scored.sort(key=lambda item: item[0], reverse=True)
        section_matches[section] = [item[1] for item in scored[:top_k]]

    return section_matches


def _rewrite_plan_with_rag_evidence(
    base_plan_markdown: str,
    section_evidence: Dict[str, List[Dict[str, Any]]],
) -> Tuple[str, List[Dict[str, str]]]:
    """Augment business plan narrative with RAG-backed evidence snippets."""

    evidence_lines: List[str] = [
        "## Evidence-backed research addendum",
        "The following evidence was retrieved from uploaded documents and integrated into the plan:",
        "",
    ]
    citations: List[Dict[str, str]] = []

    for section, chunks in section_evidence.items():
        evidence_lines.append(f"### {section}")
        if not chunks:
            evidence_lines.append("- No relevant evidence retrieved for this section.")
            evidence_lines.append("")
            continue
        for idx, chunk in enumerate(chunks, start=1):
            snippet = chunk.get("text", "").strip()[:260]
            source = chunk.get("source", "Unknown source")
            evidence_lines.append(f"- {snippet}... *(Source: {source})*")
            citations.append(
                {
                    "section": section,
                    "source": source,
                    "chunk_id": chunk.get("chunk_id", f"{source}-{idx}"),
                }
            )
        evidence_lines.append("")

    rewritten = f"{base_plan_markdown}\n\n" + "\n".join(evidence_lines).strip()
    return rewritten, citations


def _tornado_chart(
    what_if_df: pd.DataFrame, scenario_df: pd.DataFrame
) -> Optional[alt.Chart]:
    """Build a tornado chart comparing scenario deltas by NPV."""

    records: List[Dict[str, Any]] = []

    if not what_if_df.empty and {"Scenario", "NPV Δ"}.issubset(what_if_df.columns):
        for _, row in what_if_df.iterrows():
            scenario = row.get("Scenario")
            if scenario in {None, "Baseline"}:
                continue
            delta = row.get("NPV Δ")
            delta_val = float(delta) if pd.notna(delta) else None
            if delta_val is None:
                continue
            records.append(
                {
                    "Scenario": str(scenario),
                    "Impact": delta_val,
                    "Source": "What-if",
                }
            )

    if not scenario_df.empty and {"Scenario", "NPV Δ"}.issubset(scenario_df.columns):
        for _, row in scenario_df.iterrows():
            scenario = row.get("Scenario")
            if scenario in {None, "Baseline"}:
                continue
            delta = row.get("NPV Δ")
            delta_val = float(delta) if pd.notna(delta) else None
            if delta_val is None:
                continue
            records.append(
                {
                    "Scenario": str(scenario),
                    "Impact": delta_val,
                    "Source": "Scenario planning",
                }
            )

    if not records:
        return None

    data = pd.DataFrame(records)
    data = data.dropna(subset=["Impact"])
    if data.empty:
        return None

    data["Label"] = data["Source"] + " – " + data["Scenario"]
    data["abs_impact"] = data["Impact"].abs()
    order = data.sort_values("abs_impact", ascending=False)["Label"].tolist()

    bars = (
        alt.Chart(data)
        .mark_bar()
        .encode(
            y=alt.Y("Label:N", sort=order, title="Scenario"),
            x=alt.X("Impact:Q", title="NPV Δ (USD)"),
            color=alt.Color("Source:N", title="Source"),
            tooltip=[
                alt.Tooltip("Source:N"),
                alt.Tooltip("Scenario:N"),
                alt.Tooltip("Impact:Q", title="NPV Δ", format=",.0f"),
            ],
        )
    )

    zero_rule = alt.Chart(pd.DataFrame({"Impact": [0]})).mark_rule(color="black")

    return bars + zero_rule.encode(x="Impact:Q")


def _monte_carlo_distribution_charts(samples_df: pd.DataFrame) -> List[Tuple[str, alt.Chart]]:
    """Return histogram+density charts for Monte Carlo outputs."""

    charts: List[Tuple[str, alt.Chart]] = []
    if samples_df.empty:
        return charts

    numeric = samples_df.apply(_to_numeric)
    metric_map = {
        "npv": "NPV",
        "irr": "IRR",
        "min_dscr": "Minimum DSCR",
    }

    for column, label in metric_map.items():
        if column not in numeric.columns:
            continue
        series = numeric[column].dropna()
        if series.empty:
            continue
        base = pd.DataFrame({label: series})
        histogram = (
            alt.Chart(base)
            .mark_bar(opacity=0.45)
            .encode(
                x=alt.X(f"{label}:Q", bin=alt.Bin(maxbins=40), title=label),
                y=alt.Y("count():Q", title="Iterations"),
                tooltip=[alt.Tooltip("count():Q", title="Iterations")],
            )
        )
        density = (
            alt.Chart(base)
            .transform_density(label, as_=[label, "density"])
            .mark_line(color="#d62728")
            .encode(
                x=alt.X(f"{label}:Q", title=label),
                y=alt.Y("density:Q", title="Density"),
            )
        )
        charts.append((label, histogram + density))

    return charts


def _build_anomaly_points(risk_df: pd.DataFrame, trend_df: pd.DataFrame) -> pd.DataFrame:
    """Align anomaly flags with historical revenue values for plotting."""

    if risk_df.empty or trend_df.empty:
        return pd.DataFrame()

    if "Year" not in risk_df.columns or {"year", "revenue"} - set(trend_df.columns):
        return pd.DataFrame()

    revenue_lookup = (
        trend_df[["year", "revenue"]]
        .assign(year=lambda df: pd.to_numeric(df["year"], errors="coerce"))
        .dropna(subset=["year", "revenue"])
        .set_index("year")["revenue"]
    )

    points: List[Dict[str, Any]] = []
    aligned = risk_df.copy()
    aligned["Year"] = _to_numeric(aligned["Year"])
    if "Observed growth" in aligned.columns:
        aligned["Observed growth"] = _to_numeric(aligned["Observed growth"])

    for _, row in aligned.dropna(subset=["Year"]).iterrows():
        year = float(row["Year"])
        if year not in revenue_lookup.index:
            continue
        points.append(
            {
                "Year": int(year),
                "Value": float(revenue_lookup.loc[year]),
                "Flag": row.get("Flag", "Anomaly"),
                "Observed growth": row.get("Observed growth"),
            }
        )

    return pd.DataFrame(points)


def _forecast_overlay_chart(
    trend_df: pd.DataFrame,
    forecast_df: pd.DataFrame,
    time_series_df: pd.DataFrame,
    anomaly_points: pd.DataFrame,
) -> Optional[alt.Chart]:
    """Render a combined forecast chart with anomaly highlights."""

    records: List[Dict[str, Any]] = []

    if not trend_df.empty and {"year", "revenue"}.issubset(trend_df.columns):
        temp = trend_df.copy()
        temp["year"] = _to_numeric(temp["year"])
        temp["revenue"] = _to_numeric(temp["revenue"])
        for _, row in temp.dropna(subset=["year", "revenue"]).iterrows():
            records.append(
                {
                    "Year": int(row["year"]),
                    "Series": "Historical revenue",
                    "Value": float(row["revenue"]),
                }
            )

    if not forecast_df.empty and {"Year", "Revenue forecast"}.issubset(forecast_df.columns):
        temp = forecast_df.copy()
        temp["Year"] = _to_numeric(temp["Year"])
        temp["Revenue forecast"] = _to_numeric(temp["Revenue forecast"])
        for _, row in temp.dropna(subset=["Year", "Revenue forecast"]).iterrows():
            records.append(
                {
                    "Year": int(row["Year"]),
                    "Series": "Automated forecast",
                    "Value": float(row["Revenue forecast"]),
                }
            )

    if not time_series_df.empty and {"Year", "Revenue forecast"}.issubset(time_series_df.columns):
        temp = time_series_df.copy()
        temp["Year"] = _to_numeric(temp["Year"])
        temp["Revenue forecast"] = _to_numeric(temp["Revenue forecast"])
        for _, row in temp.dropna(subset=["Year", "Revenue forecast"]).iterrows():
            records.append(
                {
                    "Year": int(row["Year"]),
                    "Series": "AR(1) forecast",
                    "Value": float(row["Revenue forecast"]),
                }
            )

    if not records:
        return None

    data = pd.DataFrame(records)
    data = data.dropna(subset=["Year", "Value"])
    if data.empty:
        return None

    data["Year"] = data["Year"].astype(int)

    base_chart = (
        alt.Chart(data)
        .mark_line(point=True)
        .encode(
            x=alt.X("Year:O", title="Year"),
            y=alt.Y("Value:Q", title="Revenue (USD)"),
            color=alt.Color("Series:N", title="Series"),
            tooltip=[
                alt.Tooltip("Series:N"),
                alt.Tooltip("Year:O"),
                alt.Tooltip("Value:Q", title="Revenue", format=",.0f"),
            ],
        )
    )

    if anomaly_points.empty:
        return base_chart

    anomalies = anomaly_points.copy()
    anomalies["Year"] = _to_numeric(anomalies["Year"])
    anomalies["Value"] = _to_numeric(anomalies["Value"])
    anomalies = anomalies.dropna(subset=["Year", "Value"])
    anomalies["Year"] = anomalies["Year"].astype(int)

    points = (
        alt.Chart(anomalies)
        .mark_point(shape="triangle-up", size=120, color="#d62728")
        .encode(
            x=alt.X("Year:O"),
            y=alt.Y("Value:Q"),
            tooltip=[
                alt.Tooltip("Flag:N"),
                alt.Tooltip("Year:O"),
                alt.Tooltip("Value:Q", title="Revenue", format=",.0f"),
                alt.Tooltip("Observed growth:Q", format=".2%"),
            ],
        )
    )

    return base_chart + points


def _break_even_heatmap_charts(
    break_even_df: pd.DataFrame,
) -> Tuple[Optional[alt.Chart], Optional[alt.Chart]]:
    """Build cost heatmap and break-even unit charts for the break-even analysis."""

    if break_even_df.empty or "Category" not in break_even_df.columns:
        return None, None

    cost_columns = [col for col in ["Direct cost", "Shared cost"] if col in break_even_df.columns]
    heatmap: Optional[alt.Chart] = None
    units_chart: Optional[alt.Chart] = None

    if cost_columns:
        cost_long = (
            break_even_df.melt(
                id_vars="Category",
                value_vars=cost_columns,
                var_name="Cost component",
                value_name="Cost",
            )
        )
        cost_long["Cost"] = _to_numeric(cost_long["Cost"])
        cost_long = cost_long.dropna(subset=["Cost"])
        if not cost_long.empty:
            heatmap_base = (
                alt.Chart(cost_long)
                .mark_rect()
                .encode(
                    x=alt.X("Category:N", title="Category"),
                    y=alt.Y("Cost component:N", title="Cost type"),
                    color=alt.Color(
                        "Cost:Q",
                        title="Annual cost (USD)",
                        scale=alt.Scale(scheme="yellowgreenblue"),
                    ),
                    tooltip=[
                        alt.Tooltip("Category:N"),
                        alt.Tooltip("Cost component:N"),
                        alt.Tooltip("Cost:Q", format=",.0f"),
                    ],
                )
            )
            heatmap_text = (
                alt.Chart(cost_long)
                .mark_text(color="black")
                .encode(
                    x=alt.X("Category:N"),
                    y=alt.Y("Cost component:N"),
                    text=alt.Text("Cost:Q", format=",.0f"),
                )
            )
            heatmap = heatmap_base + heatmap_text

    if "Break-even units" in break_even_df.columns:
        units_data = break_even_df[["Category", "Break-even units"]].copy()
        units_data["Break-even units"] = _to_numeric(units_data["Break-even units"])
        units_data = units_data.dropna(subset=["Break-even units"])
        if not units_data.empty:
            units_chart = (
                alt.Chart(units_data)
                .mark_bar()
                .encode(
                    x=alt.X("Category:N", title="Category"),
                    y=alt.Y("Break-even units:Q", title="Break-even units"),
                    color=alt.Color("Category:N", legend=None),
                    tooltip=[
                        alt.Tooltip("Category:N"),
                        alt.Tooltip("Break-even units:Q", format=",.0f"),
                    ],
                )
            )

    return heatmap, units_chart

def _payload_to_ai_settings(payload: Optional[Dict[str, Any]]) -> Dict[str, Any]:
    """Return AI settings merged with defaults."""

    base = (payload or {}).get("ai_settings", {})
    settings = DEFAULT_AI_SETTINGS.copy()
    settings.update({k: v for k, v in base.items() if v is not None})
    return settings


def _ai_settings_to_payload(settings: Dict[str, Any], payload: Dict[str, Any]) -> None:
    """Persist AI settings back to the payload."""

    payload["ai_settings"] = settings


def _classify_question_intent(question: str) -> str:
    """Classify user question into a finance-analysis intent."""

    q = question.lower()
    intent_rules = [
        ("debt_capacity", ["debt", "dscr", "coverage", "leverage", "covenant"]),
        ("risk", ["risk", "volatility", "downside", "stress", "uncertain"]),
        ("scenario", ["scenario", "what-if", "sensitivity", "case", "simulate"]),
        ("operational_kpi", ["production", "cycle", "mortality", "feed", "revenue", "kpi"]),
        ("fundraising_narrative", ["investor", "fundraising", "pitch", "memo", "investment"]),
        ("valuation", ["npv", "irr", "value", "valuation", "return", "payback"]),
    ]
    for intent, keywords in intent_rules:
        if any(keyword in q for keyword in keywords):
            return intent
    return "valuation"


def _compute_chat_kpis(
    context_frames: Dict[str, pd.DataFrame], valuation: Dict[str, Any]
) -> Dict[str, float]:
    """Compute deterministic KPI snapshot from model outputs."""

    kpis: Dict[str, float] = {}
    npv = valuation.get("npv")
    irr = valuation.get("irr")
    if isinstance(npv, (int, float)):
        kpis["npv"] = float(npv)
    if isinstance(irr, (int, float)):
        kpis["irr"] = float(irr)

    dscr_df = context_frames.get("DSCR", pd.DataFrame())
    if not dscr_df.empty and "dscr" in dscr_df.columns:
        values = pd.to_numeric(dscr_df["dscr"], errors="coerce").dropna()
        if not values.empty:
            kpis["avg_dscr"] = float(values.mean())
            kpis["min_dscr"] = float(values.min())

    annual_totals_df = context_frames.get("Revenue annual totals", pd.DataFrame())
    if not annual_totals_df.empty and {"Year", "Revenue"}.issubset(annual_totals_df.columns):
        rev = annual_totals_df.copy()
        rev["Revenue"] = pd.to_numeric(rev["Revenue"], errors="coerce")
        rev = rev.dropna(subset=["Revenue"])
        if len(rev) >= 2 and rev["Revenue"].iloc[0] > 0:
            start = float(rev["Revenue"].iloc[0])
            end = float(rev["Revenue"].iloc[-1])
            years = max(len(rev) - 1, 1)
            kpis["revenue_cagr"] = (end / start) ** (1 / years) - 1

    cashflow_df = context_frames.get("Cash flows", pd.DataFrame())
    if not cashflow_df.empty and "free_cash_flow" in cashflow_df.columns:
        fcf = pd.to_numeric(cashflow_df["free_cash_flow"], errors="coerce").dropna()
        if not fcf.empty:
            kpis["latest_fcf"] = float(fcf.iloc[-1])

    return kpis


def _compute_sensitivity_impacts(kpis: Dict[str, float]) -> Dict[str, float]:
    """Produce deterministic, transparent sensitivity proxies from KPI baseline."""

    npv = kpis.get("npv", 0.0)
    return {
        "price_minus_5pct_npv": npv * 0.85,
        "price_plus_5pct_npv": npv * 1.15,
        "cost_plus_5pct_npv": npv * 0.9,
    }


def _select_evidence_rows(
    question: str, context_frames: Dict[str, pd.DataFrame], intent: str, top_k: int = 5
) -> List[Dict[str, Any]]:
    """Select top evidence rows from intent-relevant tables only."""

    intent_sources = {
        "valuation": ["Valuation", "Cash flows", "Advanced metrics"],
        "debt_capacity": ["DSCR", "Coverage analysis", "Leverage analysis", "Debt schedule"],
        "risk": ["Monte Carlo summary", "Monte Carlo samples", "Risk observations", "Scenario planning"],
        "scenario": ["What-if analysis", "Scenario planning", "Forecast projections"],
        "operational_kpi": ["Production Cycles", "Annual summary", "Revenue annual totals", "Revenue by category"],
        "fundraising_narrative": ["Valuation", "DSCR", "Revenue annual totals", "Risk observations"],
    }
    allowed = set(intent_sources.get(intent, []))
    query_tokens = set(_tokenize_text(question))
    scored: List[Tuple[int, str, int, str]] = []

    for source, frame in context_frames.items():
        if allowed and source not in allowed:
            continue
        if frame.empty:
            continue
        for idx, row in enumerate(frame.head(20).fillna("").astype(str).to_dict("records"), start=1):
            row_text = "; ".join(f"{k}: {v}" for k, v in row.items() if str(v).strip())
            score = len(query_tokens.intersection(set(_tokenize_text(row_text))))
            if score > 0:
                scored.append((score, source, idx, row_text))
    scored.sort(key=lambda x: x[0], reverse=True)
    evidence_rows: List[Dict[str, Any]] = []
    for score, source, row_idx, row_text in scored[:top_k]:
        evidence_rows.append(
            {
                "score": score,
                "source": source,
                "row_idx": row_idx,
                "row_text": row_text,
                "citation": f"[Source: {source}, Row {row_idx}]",
            }
        )
    return evidence_rows


def _compute_confidence_score(intent: str, kpis: Dict[str, float], evidence_count: int) -> float:
    """Compute a simple confidence score for response quality control."""

    required = {
        "valuation": ["npv", "irr"],
        "debt_capacity": ["avg_dscr", "min_dscr"],
        "risk": ["npv"],
        "scenario": ["npv"],
        "operational_kpi": ["latest_fcf"],
        "fundraising_narrative": ["npv", "avg_dscr"],
    }.get(intent, ["npv"])
    available_required = sum(1 for key in required if key in kpis)
    kpi_coverage = available_required / max(len(required), 1)
    evidence_score = min(evidence_count / 5.0, 1.0)
    return 0.6 * kpi_coverage + 0.4 * evidence_score


def _validate_consistency(
    kpis: Dict[str, float], sensitivity: Dict[str, float]
) -> List[str]:
    """Run basic consistency checks to suppress contradictory narratives."""

    issues: List[str] = []
    avg_dscr = kpis.get("avg_dscr")
    min_dscr = kpis.get("min_dscr")
    if avg_dscr is not None and min_dscr is not None and min_dscr > avg_dscr:
        issues.append("Minimum DSCR is greater than average DSCR; check DSCR inputs.")
    p5 = sensitivity.get("price_minus_5pct_npv")
    p95 = sensitivity.get("price_plus_5pct_npv")
    if p5 is not None and p95 is not None and p5 > p95:
        issues.append("Sensitivity ordering appears inverted (P5 greater than P95).")
    return issues


def _run_true_deterministic_sensitivity(
    assumptions: Assumptions,
) -> Dict[str, float]:
    """Run deterministic reruns with shocked assumptions for true sensitivity deltas."""

    base_results = generate_model_outputs(assumptions, analytics_plan=AnalyticsPlan.summary())
    base_npv = _to_float(base_results.get("valuation", {}).get("npv")) or 0.0

    def _npv_for(**overrides: Any) -> float:
        shocked = Assumptions(**{**asdict(assumptions), **overrides})
        out = generate_model_outputs(shocked, analytics_plan=AnalyticsPlan.summary())
        return _to_float(out.get("valuation", {}).get("npv")) or float("nan")

    price_minus = _npv_for(live_price_per_kg=float(assumptions.live_price_per_kg) * 0.95)
    price_plus = _npv_for(live_price_per_kg=float(assumptions.live_price_per_kg) * 1.05)
    cost_plus = _npv_for(feed_cost_per_kg=float(assumptions.feed_cost_per_kg) * 1.05)

    return {
        "base_npv": base_npv,
        "price_minus_5pct_npv": price_minus,
        "price_plus_5pct_npv": price_plus,
        "cost_plus_5pct_npv": cost_plus,
    }




def _retrieve_rag_evidence_for_question(
    question: str,
    chunks: Optional[List[Dict[str, Any]]],
    top_k: int = 3,
) -> List[Dict[str, Any]]:
    """Return top lexical RAG chunk matches for a user question."""

    if not chunks:
        return []

    query_tokens = set(_simple_tokenise(question))
    if not query_tokens:
        return []

    scored: List[Tuple[float, Dict[str, Any]]] = []
    for chunk in chunks:
        text = chunk.get("text", "")
        tokens = set(_simple_tokenise(text))
        if not tokens:
            continue
        overlap = len(query_tokens & tokens)
        if overlap <= 0:
            continue
        score = overlap / max(len(query_tokens), 1)
        scored.append((score, chunk))

    scored.sort(key=lambda item: item[0], reverse=True)
    return [chunk for _, chunk in scored[:top_k]]


def _answer_model_question(
    question: str,
    context_frames: Dict[str, pd.DataFrame],
    assumptions: Assumptions,
    valuation: Dict[str, Any],
    rag_chunks: Optional[List[Dict[str, Any]]] = None,
) -> str:
    """Return a direct answer grounded in model outputs plus optional RAG evidence."""

    q = (question or "").strip()
    if not q:
        return "Please enter a question about the model outputs."

    intent = _classify_question_intent(q)
    kpis = _compute_chat_kpis(context_frames, valuation)
    sensitivity = _run_true_deterministic_sensitivity(assumptions)
    evidence_rows = _select_evidence_rows(q, context_frames, intent, top_k=5)
    confidence_score = _compute_confidence_score(intent, kpis, len(evidence_rows))
    consistency_issues = _validate_consistency(kpis, sensitivity)

    if confidence_score < 0.35:
        return (
            "I don’t have enough high-confidence evidence from the current model outputs to answer reliably. "
            "Please run the relevant simulation blocks (Monte Carlo / break-even / scenario planning) and ask again."
        )
    if consistency_issues:
        issues_text = " ".join(consistency_issues)
        return (
            "I detected a model consistency issue and will avoid a potentially misleading conclusion. "
            f"Please review inputs/results first: {issues_text}"
        )

    rag_matches = _retrieve_rag_evidence_for_question(q, rag_chunks, top_k=3)

    parts = [
        "Direct answer from current model outputs:",
        f"- NPV: {kpis.get('npv', float('nan')):,.0f}" if "npv" in kpis else "- NPV: N/A in current outputs",
        f"- IRR: {kpis.get('irr', float('nan')):.2%}" if "irr" in kpis else "- IRR: N/A in current outputs",
        f"- Avg DSCR: {kpis.get('avg_dscr', float('nan')):.2f}" if "avg_dscr" in kpis else "- Avg DSCR: N/A in current outputs",
        f"- Min DSCR: {kpis.get('min_dscr', float('nan')):.2f}" if "min_dscr" in kpis else "- Min DSCR: N/A in current outputs",
        f"- Latest FCF: {kpis.get('latest_fcf', float('nan')):,.0f}" if "latest_fcf" in kpis else "- Latest FCF: N/A in current outputs",
        f"- NPV if price -5%: {sensitivity['price_minus_5pct_npv']:,.0f}",
        f"- NPV if price +5%: {sensitivity['price_plus_5pct_npv']:,.0f}",
        f"- NPV if costs +5%: {sensitivity['cost_plus_5pct_npv']:,.0f}",
        "",
        "I used model tables/results in this scenario and (when available) indexed RAG snippets from uploaded files.",
        f"Matched intent: {intent.replace('_', ' ').title()} | Confidence score: {confidence_score:.2f}",
    ]

    if evidence_rows:
        parts.append("### Evidence mapping (model tables)")
        for item in evidence_rows:
            parts.append(f"- {item['citation']} {item['row_text']}")

    if rag_matches:
        parts.append("### Evidence mapping (RAG documents)")
        for match in rag_matches:
            snippet = " ".join(str(match.get("text", "")).split())[:280]
            parts.append(
                f"- [{match.get('source', 'document')}] {snippet}... (chunk: {match.get('chunk_id', 'n/a')})"
            )

    return "\n".join(parts)


def _assess_investor_readiness(
    kpis: Dict[str, float], benchmarks: Dict[str, float]
) -> Dict[str, Any]:
    """Score investor readiness against benchmark hurdles with explainable outcomes."""

    checks: List[Dict[str, Any]] = []
    irr_value = kpis.get("irr")
    dscr_value = kpis.get("avg_dscr")
    payback_value = kpis.get("payback_years")

    target_irr = float(benchmarks.get("target_irr", 0.18))
    min_dscr = float(benchmarks.get("min_dscr", 1.5))
    max_payback = float(benchmarks.get("max_payback_years", 6.0))

    checks.append(
        {
            "metric": "IRR",
            "value": irr_value,
            "target": target_irr,
            "pass": irr_value is not None and irr_value >= target_irr,
            "rule": "Higher is better",
        }
    )
    checks.append(
        {
            "metric": "Average DSCR",
            "value": dscr_value,
            "target": min_dscr,
            "pass": dscr_value is not None and dscr_value >= min_dscr,
            "rule": "Higher is better",
        }
    )
    checks.append(
        {
            "metric": "Payback (years)",
            "value": payback_value,
            "target": max_payback,
            "pass": payback_value is not None and payback_value <= max_payback,
            "rule": "Lower is better",
        }
    )

    passes = sum(1 for item in checks if item["pass"])
    score = (passes / max(len(checks), 1)) * 100.0
    if score >= 80:
        verdict = "Investor-ready"
    elif score >= 50:
        verdict = "Near-ready"
    else:
        verdict = "Needs improvement"
    return {"score": score, "verdict": verdict, "checks": checks}


def _build_unified_orchestration_state(
    config: Dict[str, Any],
    context_frames: Dict[str, pd.DataFrame],
    assumptions: Assumptions,
    valuation: Dict[str, Any],
    rag_chunks: Optional[List[Dict[str, Any]]],
    benchmarks: Dict[str, float],
    investor_recommendations: List[str],
) -> Dict[str, Any]:
    """Create shared knowledge + reasoning outputs for the unified AI orchestration layer."""

    kpis = _compute_chat_kpis(context_frames, valuation)
    payback_years = _to_float(config.get("payback_years"))
    if payback_years is not None:
        kpis["payback_years"] = payback_years
    sensitivity = _run_true_deterministic_sensitivity(assumptions)
    readiness = _assess_investor_readiness(kpis, benchmarks)

    strategic_analysis: List[str] = []
    if "avg_dscr" in kpis and kpis["avg_dscr"] < benchmarks.get("min_dscr", 1.5):
        strategic_analysis.append(
            "Debt service resilience is below benchmark; prioritize margin protection and debt profile optimization."
        )
    if "irr" in kpis and kpis["irr"] < benchmarks.get("target_irr", 0.18):
        strategic_analysis.append(
            "Return profile is below investor hurdle; focus on pricing, productivity, and capex efficiency levers."
        )
    if sensitivity.get("price_minus_5pct_npv", 0.0) < 0:
        strategic_analysis.append(
            "Downside pricing stress drives negative NPV; add hedging/commercial contracts to improve downside protection."
        )
    if not strategic_analysis:
        strategic_analysis.append(
            "Current scenario shows balanced return and resilience against benchmark thresholds."
        )

    return {
        "config": config,
        "knowledge": {
            "context_frames": context_frames,
            "rag_chunks": rag_chunks or [],
            "kpis": kpis,
            "sensitivity": sensitivity,
            "benchmarks": benchmarks,
        },
        "reasoning": {
            "investor_readiness": readiness,
            "strategic_analysis": strategic_analysis,
            "proactive_recommendations": investor_recommendations[:5],
        },
    }


def _answer_unified_orchestrator_question(
    question: str,
    orchestration_state: Dict[str, Any],
) -> str:
    """Answer with shared orchestration knowledge and explainable decision-support context."""

    q = (question or "").strip()
    if not q:
        return "Please enter a question for the unified AI orchestration system."

    knowledge = orchestration_state.get("knowledge", {})
    context_frames = knowledge.get("context_frames", {})
    kpis = knowledge.get("kpis", {})
    sensitivity = knowledge.get("sensitivity", {})
    readiness = orchestration_state.get("reasoning", {}).get("investor_readiness", {})
    strategic_analysis = orchestration_state.get("reasoning", {}).get(
        "strategic_analysis", []
    )
    recommendations = orchestration_state.get("reasoning", {}).get(
        "proactive_recommendations", []
    )

    intent = _classify_question_intent(q)
    evidence_rows = _select_evidence_rows(q, context_frames, intent, top_k=5)
    rag_matches = _retrieve_rag_evidence_for_question(
        q, knowledge.get("rag_chunks", []), top_k=3
    )

    lines = [
        "Unified AI orchestration response:",
        f"- Intent: {intent.replace('_', ' ').title()}",
        f"- Investor-readiness: {readiness.get('verdict', 'N/A')} ({readiness.get('score', 0.0):.0f}/100)",
        f"- NPV: {kpis.get('npv', float('nan')):,.0f}" if "npv" in kpis else "- NPV: N/A",
        f"- IRR: {kpis.get('irr', float('nan')):.2%}" if "irr" in kpis else "- IRR: N/A",
        f"- Avg DSCR: {kpis.get('avg_dscr', float('nan')):.2f}" if "avg_dscr" in kpis else "- Avg DSCR: N/A",
        f"- Price -5% NPV: {sensitivity.get('price_minus_5pct_npv', float('nan')):,.0f}",
    ]
    if strategic_analysis:
        lines.append("")
        lines.append("Strategic analysis:")
        lines.extend([f"- {item}" for item in strategic_analysis[:3]])
    if recommendations:
        lines.append("")
        lines.append("Proactive recommendations:")
        lines.extend([f"- {item}" for item in recommendations[:3]])
    if evidence_rows:
        lines.append("")
        lines.append("Explainability (model evidence):")
        for item in evidence_rows:
            lines.append(f"- {item['citation']} {item['row_text']}")
    if rag_matches:
        lines.append("")
        lines.append("Grounding (research evidence):")
        for match in rag_matches:
            snippet = " ".join(str(match.get("text", "")).split())[:220]
            source = match.get("source", "document")
            lines.append(f"- [RAG: {source}] {snippet}")
    return "\n".join(lines)


def _rerun() -> None:
    """Trigger a Streamlit rerun."""

    if hasattr(st, "rerun"):
        try:
            st.rerun()
            return
        except Exception:
            pass
    # Fallback for Streamlit versions that still expose experimental API.
    st.experimental_rerun()


def _ensure_scenario_payload(
    selected_scenario: str, snapshot: Dict[str, Any]
) -> Tuple[ScenarioModel, Dict[str, Any]]:
    """Ensure the scenario model/results are cached for the provided snapshot."""

    cache = st.session_state.setdefault("scenario_payloads", {})
    existing = cache.get(selected_scenario)
    if existing and existing.get("snapshot") == snapshot:
        return existing["model"], existing["results"]

    assumptions_data = snapshot.get("assumptions", {})
    assumptions = Assumptions(**assumptions_data) if assumptions_data else Assumptions()
    model = ScenarioModel(scenario=selected_scenario, assumptions=assumptions)
    results = generate_model_outputs(
        assumptions,
        analytics_plan=AnalyticsPlan.summary(),
    )
    cache[selected_scenario] = {
        "snapshot": copy.deepcopy(snapshot),
        "model": model,
        "results": results,
    }
    return model, results


def _normalize_scenario_payload(payload: Any) -> Dict[str, Any]:
    """Normalize persisted scenario payloads from older/newer app versions."""

    normalized: Dict[str, Any] = payload if isinstance(payload, dict) else {}
    assumptions_payload = normalized.get("assumptions", {})
    ai_payload = normalized.get("ai_settings", {})
    benchmark_payload = normalized.get("investor_benchmarks", {})

    if not isinstance(assumptions_payload, dict):
        assumptions_payload = {}
    if not isinstance(ai_payload, dict):
        ai_payload = {}
    if not isinstance(benchmark_payload, dict):
        benchmark_payload = {}

    if "eggs_per_cycle_default" in assumptions_payload and "eggs_per_bird_per_cycle" not in assumptions_payload:
        assumptions_payload["eggs_per_bird_per_cycle"] = float(
            assumptions_payload.get("eggs_per_cycle_default", 0.0) or 0.0
        ) / max(float(assumptions_payload.get("birds_per_cycle", 1) or 1), 1.0)
    assumptions_payload.pop("eggs_per_cycle_default", None)

    normalized["assumptions"] = assumptions_payload
    normalized["ai_settings"] = {
        **DEFAULT_AI_SETTINGS,
        **{k: v for k, v in ai_payload.items() if v is not None},
    }
    normalized["investor_benchmarks"] = {
        **DEFAULT_INVESTOR_BENCHMARKS,
        **{k: v for k, v in benchmark_payload.items() if v is not None},
    }
    return normalized


def _generate_excel_bytes(
    model: ScenarioModel, results: Dict[str, Any], scenario: str
) -> bytes:
    """Create an Excel workbook in memory for the supplied scenario results."""

    buffer = BytesIO()
    try:
        import xlsxwriter  # type: ignore  # noqa: F401

        engine = "xlsxwriter"
    except ImportError:
        try:
            import openpyxl  # type: ignore  # noqa: F401

            engine = "openpyxl"
        except ImportError as exc:  # pragma: no cover - surfaced in UI
            st.error(
                "Excel exports require either the `xlsxwriter` or `openpyxl` package. "
                "Install one of these dependencies and try again."
            )
            raise RuntimeError("Missing Excel writer dependency") from exc

    def _excel_ready_df(frame: pd.DataFrame) -> pd.DataFrame:
        safe = frame.copy()
        if safe.empty:
            return safe
        for column in safe.columns:
            if safe[column].dtype == object:
                safe[column] = safe[column].apply(
                    lambda value: json.dumps(value)
                    if isinstance(value, (dict, list, tuple, set))
                    else value
                )
        return safe

    try:
        with pd.ExcelWriter(buffer, engine=engine) as writer:
            _excel_ready_df(pd.DataFrame(results["assumptions_schedule"])).to_excel(
                writer, sheet_name="Assumptions", index=False
            )
            _excel_ready_df(pd.DataFrame([asdict(model.assumptions)])).to_excel(
                writer, sheet_name="Input Values", index=False
            )
            _excel_ready_df(pd.DataFrame([asdict(cycle) for cycle in results["cycles"]])).to_excel(
                writer, sheet_name="Production Cycles", index=False
            )
            _excel_ready_df(pd.DataFrame([asdict(results["annual"])])).to_excel(
                writer, sheet_name="Annual Summary", index=False
            )
            _excel_ready_df(pd.DataFrame([asdict(row) for row in results["cashflows"]])).to_excel(
                writer, sheet_name="Cash Flows", index=False
            )
            _excel_ready_df(pd.DataFrame([results["valuation"]])).to_excel(
                writer, sheet_name="Valuation", index=False
            )

            financials = results["financial_statements"]
            _excel_ready_df(pd.DataFrame([asdict(row) for row in financials["income_statement"]])).to_excel(
                writer, sheet_name="Income Statement", index=False
            )
            _excel_ready_df(pd.DataFrame([asdict(row) for row in financials["balance_sheet"]])).to_excel(
                writer, sheet_name="Balance Sheet", index=False
            )
            _excel_ready_df(pd.DataFrame([asdict(row) for row in financials["cash_flow_statement"]])).to_excel(
                writer, sheet_name="Cash Flow Statement", index=False
            )
            _excel_ready_df(pd.DataFrame(financials["loan_schedule"])).to_excel(
                writer, sheet_name="Debt Schedule", index=False
            )

            advanced = results["advanced_analytics"]
            _excel_ready_df(pd.DataFrame(advanced["metrics"])).to_excel(
                writer, sheet_name="Advanced Metrics", index=False
            )
            _excel_ready_df(pd.DataFrame(advanced["dscr"])).to_excel(
                writer, sheet_name="DSCR", index=False
            )
            _excel_ready_df(pd.DataFrame(advanced["trend"])).to_excel(
                writer, sheet_name="Trend Analysis", index=False
            )

            for category, rows in results["revenue_schedules"].items():
                safe_name = (
                    category.replace("(", "")
                    .replace(")", "")
                    .replace(",", "")
                    .replace("-", " ")
                )
                safe_name = " ".join(word.title() for word in safe_name.split())
                _excel_ready_df(pd.DataFrame(rows)).to_excel(
                    writer, sheet_name=safe_name[:31] or "Revenue", index=False
                )

            ai_settings = st.session_state.get("ai_settings", DEFAULT_AI_SETTINGS)
            _excel_ready_df(pd.DataFrame([ai_settings])).to_excel(
                writer, sheet_name="AI Settings", index=False
            )

            workbook = writer.book
            if engine == "xlsxwriter":
                workbook.set_properties(
                    {
                        "title": f"Broiler Model - {scenario}",
                        "subject": "Broiler chicken financial model",
                        "comments": "Generated via Streamlit dashboard",
                    }
                )
            else:
                props = workbook.properties
                props.title = f"Broiler Model - {scenario}"
                props.subject = "Broiler chicken financial model"
                props.comments = "Generated via Streamlit dashboard"
    except Exception as exc:  # pragma: no cover - surfaced in UI
        st.error(f"Excel export failed due to serialization issue: {exc}")
        raise RuntimeError("Excel export serialization failure") from exc

    buffer.seek(0)
    return buffer.getvalue()


def _excel_dependency_health() -> Dict[str, Any]:
    """Return availability state for optional Excel writer dependencies."""

    status = {"xlsxwriter": False, "openpyxl": False}
    try:
        import xlsxwriter  # type: ignore  # noqa: F401

        status["xlsxwriter"] = True
    except ImportError:
        pass
    try:
        import openpyxl  # type: ignore  # noqa: F401

        status["openpyxl"] = True
    except ImportError:
        pass
    status["ready"] = bool(status["xlsxwriter"] or status["openpyxl"])
    return status


def _generate_csv_zip_bytes(
    model: ScenarioModel, results: Dict[str, Any], scenario: str
) -> bytes:
    """Create a ZIP bundle of CSV reports in memory for the supplied scenario results."""

    def _csv_ready_df(frame: pd.DataFrame) -> pd.DataFrame:
        safe = frame.copy()
        if safe.empty:
            return safe
        for column in safe.columns:
            if safe[column].dtype == object:
                safe[column] = safe[column].apply(
                    lambda value: json.dumps(value)
                    if isinstance(value, (dict, list, tuple, set))
                    else value
                )
        return safe

    export_frames: Dict[str, pd.DataFrame] = {
        "assumptions_schedule.csv": pd.DataFrame(results["assumptions_schedule"]),
        "input_values.csv": pd.DataFrame([asdict(model.assumptions)]),
        "production_cycles.csv": pd.DataFrame([asdict(cycle) for cycle in results["cycles"]]),
        "annual_summary.csv": pd.DataFrame([asdict(results["annual"])]),
        "cash_flows.csv": pd.DataFrame([asdict(row) for row in results["cashflows"]]),
        "valuation.csv": pd.DataFrame([results["valuation"]]),
    }

    financials = results["financial_statements"]
    export_frames.update(
        {
            "income_statement.csv": pd.DataFrame(
                [asdict(row) for row in financials["income_statement"]]
            ),
            "balance_sheet.csv": pd.DataFrame(
                [asdict(row) for row in financials["balance_sheet"]]
            ),
            "cash_flow_statement.csv": pd.DataFrame(
                [asdict(row) for row in financials["cash_flow_statement"]]
            ),
            "debt_schedule.csv": pd.DataFrame(financials["loan_schedule"]),
        }
    )

    advanced = results["advanced_analytics"]
    export_frames.update(
        {
            "advanced_metrics.csv": pd.DataFrame(advanced["metrics"]),
            "dscr.csv": pd.DataFrame(advanced["dscr"]),
            "trend_analysis.csv": pd.DataFrame(advanced["trend"]),
        }
    )

    for category, rows in results["revenue_schedules"].items():
        safe_name = re.sub(r"[^a-z0-9]+", "_", category.lower()).strip("_")
        export_frames[f"revenue_{safe_name or 'schedule'}.csv"] = pd.DataFrame(rows)

    ai_settings = st.session_state.get("ai_settings", DEFAULT_AI_SETTINGS)
    export_frames["ai_settings.csv"] = pd.DataFrame([ai_settings])

    buffer = BytesIO()
    with zipfile.ZipFile(buffer, mode="w", compression=zipfile.ZIP_DEFLATED) as archive:
        archive.writestr(
            "README.txt",
            (
                f"Broiler chicken financial model export for scenario: {scenario}\n"
                "This archive contains CSV versions of report tabs that are normally "
                "included in the Excel workbook export.\n"
            ),
        )
        for filename, frame in export_frames.items():
            archive.writestr(filename, _csv_ready_df(frame).to_csv(index=False))

    buffer.seek(0)
    return buffer.getvalue()


def _generate_excel_xml_bytes(
    model: ScenarioModel, results: Dict[str, Any], scenario: str
) -> bytes:
    """Create an Excel-compatible SpreadsheetML workbook without external engines."""

    def _xml_ready_df(frame: pd.DataFrame) -> pd.DataFrame:
        safe = frame.copy()
        if safe.empty:
            return safe
        for column in safe.columns:
            if safe[column].dtype == object:
                safe[column] = safe[column].apply(
                    lambda value: json.dumps(value)
                    if isinstance(value, (dict, list, tuple, set))
                    else value
                )
        return safe

    sheets: Dict[str, pd.DataFrame] = {
        "Assumptions": pd.DataFrame(results["assumptions_schedule"]),
        "Input Values": pd.DataFrame([asdict(model.assumptions)]),
        "Production Cycles": pd.DataFrame([asdict(cycle) for cycle in results["cycles"]]),
        "Annual Summary": pd.DataFrame([asdict(results["annual"])]),
        "Cash Flows": pd.DataFrame([asdict(row) for row in results["cashflows"]]),
        "Valuation": pd.DataFrame([results["valuation"]]),
    }
    financials = results["financial_statements"]
    sheets["Income Statement"] = pd.DataFrame(
        [asdict(row) for row in financials["income_statement"]]
    )
    sheets["Balance Sheet"] = pd.DataFrame(
        [asdict(row) for row in financials["balance_sheet"]]
    )
    sheets["Cash Flow Statement"] = pd.DataFrame(
        [asdict(row) for row in financials["cash_flow_statement"]]
    )
    sheets["Debt Schedule"] = pd.DataFrame(financials["loan_schedule"])

    lines = [
        '<?xml version="1.0"?>',
        '<?mso-application progid="Excel.Sheet"?>',
        '<Workbook xmlns="urn:schemas-microsoft-com:office:spreadsheet" '
        'xmlns:o="urn:schemas-microsoft-com:office:office" '
        'xmlns:x="urn:schemas-microsoft-com:office:excel" '
        'xmlns:ss="urn:schemas-microsoft-com:office:spreadsheet">',
        "<DocumentProperties xmlns=\"urn:schemas-microsoft-com:office:office\">",
        f"<Title>{xml_escape(f'Broiler Model - {scenario}')}</Title>",
        "</DocumentProperties>",
    ]

    for raw_name, frame in sheets.items():
        sheet_name = raw_name[:31] or "Sheet"
        export_df = _xml_ready_df(frame)
        lines.append(f'<Worksheet ss:Name="{xml_escape(sheet_name)}"><Table>')
        if not export_df.empty:
            lines.append("<Row>")
            for col in export_df.columns:
                lines.append(
                    f'<Cell><Data ss:Type="String">{xml_escape(str(col))}</Data></Cell>'
                )
            lines.append("</Row>")
            for row in export_df.itertuples(index=False):
                lines.append("<Row>")
                for value in row:
                    if value is None or (isinstance(value, float) and math.isnan(value)):
                        cell_value = ""
                        cell_type = "String"
                    elif isinstance(value, (int, float)) and not isinstance(value, bool):
                        cell_value = str(value)
                        cell_type = "Number"
                    else:
                        cell_value = str(value)
                        cell_type = "String"
                    lines.append(
                        f'<Cell><Data ss:Type="{cell_type}">{xml_escape(cell_value)}</Data></Cell>'
                    )
                lines.append("</Row>")
        lines.append("</Table></Worksheet>")

    lines.append("</Workbook>")
    return "\n".join(lines).encode("utf-8")


AI_PROVIDER_OPTIONS = (
    "OpenAI",
    "Anthropic",
    "Azure OpenAI",
    "Vertex AI",
    "Custom",
)

ML_METHOD_LABELS: Dict[str, str] = {
    "linear_regression": "Linear Regression",
    "exponential_smoothing": "Exponential Smoothing",
    "random_forest": "Random Forest",
    "prophet": "Prophet Forecast",
}

ML_LABEL_TO_CODE = {label: code for code, label in ML_METHOD_LABELS.items()}

GEN_AI_FEATURE_LABELS: Dict[str, str] = {
    "summary": "Executive Summary",
    "risk_review": "Risk Review",
    "scenario_comparison": "Scenario Comparison",
    "cash_flow_focus": "Cash Flow Focus",
}

GEN_AI_LABEL_TO_CODE = {label: code for code, label in GEN_AI_FEATURE_LABELS.items()}

SCENARIO_OPTIONS = ["Baseline", "Expansion", "Downside"]

DEFAULT_AI_SETTINGS: Dict[str, Any] = {
    "enabled": False,
    "provider": "OpenAI",
    "model": "gpt-4o-mini",
    "forecast_horizon": 3,
    "auto_run_predictive": False,
    "ml_methods": ["linear_regression"],
    "generative_features": ["summary"],
    "api_key": "",
}

DEFAULT_INVESTOR_BENCHMARKS: Dict[str, float] = {
    "target_irr": 0.18,
    "min_dscr": 1.5,
    "max_payback_years": 6.0,
}


@dataclass
class ScenarioModel:
    """Lightweight wrapper storing the active scenario and assumptions."""

    scenario: str
    assumptions: Assumptions


def _normalise_schedule_rows(rows: List[Any]) -> List[Dict[str, Any]]:
    """Return a list of dictionaries for any supported row payload."""

    normalised: List[Dict[str, Any]] = []
    for row in rows:
        if row is None:
            continue
        if isinstance(row, dict):
            normalised.append(copy.deepcopy(row))
        elif is_dataclass(row):
            normalised.append(asdict(row))
        else:
            try:
                normalised.append(dict(row))
            except TypeError:
                raise TypeError(
                    "Schedule defaults must be dict-like or dataclass instances."
                ) from None
    return normalised


def _initialise_schedule_state(
    namespace: str,
    scenario: str,
    schedule_key: str,
    default_rows: List[Dict[str, Any]],
) -> Tuple[pd.DataFrame, Dict[str, Any]]:
    """Return the stored DataFrame for a schedule, initialising from defaults as needed."""

    store = st.session_state.setdefault(namespace, {})
    scenario_store = store.setdefault(scenario, {})
    default_records = _normalise_schedule_rows(default_rows)
    state = scenario_store.get(schedule_key)
    if state is None:
        state = {
            "data": copy.deepcopy(default_records),
            "default": copy.deepcopy(default_records),
        }
        scenario_store[schedule_key] = state
    else:
        stored_data = state.get("data", [])
        stored_default = state.get("default", [])
        if stored_data == stored_default and stored_default != default_records:
            state["data"] = copy.deepcopy(default_records)
            state["default"] = copy.deepcopy(default_records)
        else:
            if default_records and stored_data:
                sample_default = default_records[0]
                sample_stored = stored_data[0] if isinstance(stored_data[0], dict) else {}
                key_fields = [
                    field
                    for field in ("Category", "Period")
                    if field in sample_default and field in sample_stored
                ]
                if key_fields:
                    default_key_set = {
                        tuple(row.get(field) for field in key_fields)
                        for row in default_records
                    }
                    stored_key_map = {
                        tuple(row.get(field) for field in key_fields): row
                        for row in stored_data
                        if isinstance(row, dict)
                    }
                    needs_reconcile = set(stored_key_map) != default_key_set
                    if needs_reconcile:
                        reconciled_rows: List[Dict[str, Any]] = []
                        for default_row in default_records:
                            key = tuple(default_row.get(field) for field in key_fields)
                            if key in stored_key_map:
                                merged = copy.deepcopy(default_row)
                                merged.update(stored_key_map[key])
                                reconciled_rows.append(merged)
                            else:
                                reconciled_rows.append(copy.deepcopy(default_row))
                        state["data"] = reconciled_rows
                        state["default"] = copy.deepcopy(default_records)
    df = pd.DataFrame(state.get("data", default_records))
    return df, state


def _reset_scenario_edit_states(scenario: str) -> None:
    """Clear per-scenario editable table caches so new assumptions fully propagate."""

    for namespace in (
        "assumption_schedule_state",
        "revenue_schedule_state",
        "advanced_schedule_state",
        "custom_simulation_state",
    ):
        store = st.session_state.get(namespace)
        if isinstance(store, dict) and scenario in store:
            store.pop(scenario, None)
            st.session_state[namespace] = store


def _coerce_assumption_value(key: str, raw_value: Any, current_value: Any) -> Any:
    """Convert schedule entries back to the dataclass field types."""

    if raw_value is None:
        return current_value
    if pd.isna(raw_value):  # type: ignore[arg-type]
        return current_value
    if isinstance(raw_value, str):
        stripped = raw_value.strip()
        if stripped == "":
            return current_value
        raw_value = stripped

    target_type = _ASSUMPTION_FIELD_TYPES.get(key, type(current_value))

    try:
        if target_type is int:
            return int(float(raw_value))
        if target_type is float:
            return float(raw_value)
        return str(raw_value)
    except (TypeError, ValueError):
        return current_value


def _assumptions_from_schedule(
    rows: List[Dict[str, Any]],
    current: Assumptions,
) -> Optional[Assumptions]:
    """Derive an updated ``Assumptions`` instance from schedule edits."""

    updated_values = asdict(current)
    changed = False

    for row in rows:
        item = row.get("item")
        if not item:
            continue
        key = _ASSUMPTION_LABEL_TO_KEY.get(str(item))
        if not key:
            continue

        existing = updated_values.get(key)
        new_value = _coerce_assumption_value(key, row.get("value"), existing)

        if isinstance(existing, float) and isinstance(new_value, float):
            if math.isnan(existing) and math.isnan(new_value):
                continue
            if math.isclose(existing, new_value, rel_tol=1e-9, abs_tol=1e-9):
                continue
        elif existing == new_value:
            continue

        updated_values[key] = new_value
        changed = True

    if not changed:
        return None

    return Assumptions(**updated_values)


def _next_period_label(df: pd.DataFrame) -> str:
    """Generate a reasonable period label for a newly added row."""

    if "Period" not in df.columns or df["Period"].dropna().empty:
        return f"Period {len(df) + 1}"
    last_value = str(df["Period"].dropna().iloc[-1]).strip()
    parts = last_value.split()
    if parts and parts[-1].isdigit():
        prefix = " ".join(parts[:-1]) or "Period"
        return f"{prefix} {int(parts[-1]) + 1}"
    match = re.search(r"(\d+)(?!.*\d)", last_value)
    if match:
        number = int(match.group(1)) + 1
        prefix = last_value[: match.start()].strip() or "Period"
        suffix = last_value[match.end() :].strip()
        if suffix:
            return f"{prefix} {number} {suffix}".strip()
        return f"{prefix} {number}".strip()
    return f"Period {len(df) + 1}"


def _ensure_fixed_columns(
    df: pd.DataFrame, fixed_columns: Optional[Dict[str, Any]]
) -> pd.DataFrame:
    if not fixed_columns:
        return df
    new_df = df.copy()
    for column, value in fixed_columns.items():
        if column in new_df.columns:
            new_df[column] = value
    return new_df


def _add_schedule_row(
    df: pd.DataFrame,
    row_defaults: Optional[Dict[str, Any]] = None,
    fixed_columns: Optional[Dict[str, Any]] = None,
) -> pd.DataFrame:
    """Append a new row using defaults and fixed column values."""

    new_row = {col: None for col in df.columns}
    defaults = row_defaults or {}
    for column, value in defaults.items():
        if column in new_row:
            new_row[column] = copy.deepcopy(value)
    if "Period" in df.columns and "Period" not in defaults:
        new_row["Period"] = _next_period_label(df)
    if fixed_columns:
        for column, value in fixed_columns.items():
            if column in new_row:
                new_row[column] = value
    if ROW_REMOVAL_COLUMN in new_row:
        new_row[ROW_REMOVAL_COLUMN] = False
    if ROW_EDIT_COLUMN in new_row:
        new_row[ROW_EDIT_COLUMN] = False
    return pd.concat([df, pd.DataFrame([new_row])], ignore_index=True)


def _apply_yearly_increment(
    df: pd.DataFrame, columns: List[str], rate: float
) -> pd.DataFrame:
    """Apply a compound yearly increment across rows for selected columns."""

    if not columns or rate == 0.0 or df.empty:
        return df
    new_df = df.copy()
    for column in columns:
        if column not in new_df.columns:
            continue
        numeric_series = pd.to_numeric(new_df[column], errors="coerce")
        if numeric_series.notna().sum() == 0:
            # nothing numeric to increment; leave column as-is
            continue
        # ensure the column can hold floating point results from the increment
        new_df[column] = numeric_series.astype("Float64")
        prev_value: Optional[float] = None
        for idx in new_df.index:
            current = new_df.at[idx, column]
            if prev_value is None:
                if not pd.isna(current):
                    prev_value = float(current)
                continue
            prev_value = prev_value * (1 + rate)
            new_df.at[idx, column] = prev_value
    return new_df


def _auto_compute_revenue(df: pd.DataFrame) -> pd.DataFrame:
    """Recompute revenue when units and price are available."""

    if not {"Units", "Unit price", "Revenue"}.issubset(df.columns):
        return df
    new_df = df.copy()
    for idx in range(len(new_df)):
        units = pd.to_numeric(new_df.at[idx, "Units"], errors="coerce")
        price = pd.to_numeric(new_df.at[idx, "Unit price"], errors="coerce")
        if pd.isna(units) or pd.isna(price):
            continue
        new_df.at[idx, "Revenue"] = float(units) * float(price)
    return new_df


def _sanitize_key(label: str) -> str:
    return re.sub(r"[^0-9a-zA-Z_]+", "_", label.lower()).strip("_") or "schedule"


def _stringify_iterable_value(value: Any) -> Any:
    if isinstance(value, (list, tuple, dict)):
        try:
            return json.dumps(value)
        except TypeError:
            return str(value)
    return value


def _stringify_iterable_columns(df: pd.DataFrame, columns: Iterable[str]) -> pd.DataFrame:
    new_df = df.copy()
    for column in columns:
        if column in new_df.columns:
            new_df[column] = new_df[column].apply(_stringify_iterable_value)
    return new_df


def _parse_iterable_value(value: Any) -> Any:
    if value is None or value is pd.NA:
        return None
    if isinstance(value, float) and math.isnan(value):
        return None
    if isinstance(value, str):
        text = value.strip()
        if not text:
            return None
        try:
            parsed = json.loads(text)
        except json.JSONDecodeError:
            parts = [part.strip() for part in re.split(r"[,;]", text) if part.strip()]
            if not parts:
                return None
            parsed_parts: List[Any] = []
            for item in parts:
                try:
                    parsed_parts.append(float(item))
                except ValueError:
                    parsed_parts.append(item)
            parsed = parsed_parts
        return parsed
    return value


def _coerce_editor_input(value: str) -> Any:
    if value is None:
        return None
    text = str(value).strip()
    if text == "":
        return None
    try:
        if re.fullmatch(r"[-+]?\d+", text):
            return int(text)
        return float(text)
    except ValueError:
        return text


def _render_row_edit_form(
    df: pd.DataFrame,
    idx: int,
    fixed_columns: Optional[Dict[str, Any]],
    namespace: str,
    schedule_key: str,
    scenario: str,
) -> Tuple[pd.DataFrame, bool]:
    fixed_columns = fixed_columns or {}
    editable_columns = [
        column
        for column in df.columns
        if column not in {ROW_REMOVAL_COLUMN, ROW_EDIT_COLUMN}
        and column not in fixed_columns
    ]
    if not editable_columns:
        return df, False

    form_key = f"edit_row_form_{namespace}_{schedule_key}_{scenario}_{idx}"
    with st.form(form_key):
        st.markdown(f"**Editing row {idx + 1}**")
        inputs: Dict[str, str] = {}
        for column in editable_columns:
            current_value = df.at[idx, column]
            default_text = "" if pd.isna(current_value) else str(current_value)
            inputs[column] = st.text_input(column, value=default_text)
        submitted = st.form_submit_button("Save row")

    if not submitted:
        return df, False

    new_df = df.copy()
    for column, raw_value in inputs.items():
        series = new_df[column]
        if pd.api.types.is_string_dtype(series):
            new_df[column] = series.astype(object)
        new_df.at[idx, column] = _coerce_editor_input(raw_value)
    st.success(f"Row {idx + 1} updated.")
    return new_df, True


def _parse_iterable_columns(df: pd.DataFrame, columns: Iterable[str]) -> pd.DataFrame:
    new_df = df.copy()
    for column in columns:
        if column in new_df.columns:
            new_df[column] = new_df[column].apply(_parse_iterable_value)
    return new_df


def _chunked(iterable: Iterable[Any], size: int) -> Iterable[List[Any]]:
    chunk: List[Any] = []
    for item in iterable:
        chunk.append(item)
        if len(chunk) == size:
            yield chunk
            chunk = []
    if chunk:
        yield chunk


def _build_column_config(df: pd.DataFrame, *, disabled: bool, fixed: Optional[Dict[str, Any]] = None) -> Dict[str, Any]:
    """Create Streamlit column configuration based on series dtypes."""

    fixed = fixed or {}
    config: Dict[str, Any] = {}
    for column in df.columns:
        col_disabled = disabled or column in fixed
        series = df[column]
        if pd.api.types.is_bool_dtype(series):
            config[column] = st.column_config.CheckboxColumn(disabled=col_disabled)
        elif pd.api.types.is_numeric_dtype(series):
            config[column] = st.column_config.NumberColumn(disabled=col_disabled)
        elif pd.api.types.is_datetime64_any_dtype(series):
            config[column] = st.column_config.DatetimeColumn(disabled=col_disabled)
        else:
            config[column] = st.column_config.TextColumn(disabled=col_disabled)
    return config


def _render_input_section(
    title: str,
    fields: List[Dict[str, Any]],
    defaults: Assumptions,
    values: Dict[str, Any],
    *,
    columns: int = 3,
) -> None:
    """Render a consistently styled multi-column assumptions section."""

    st.markdown(f"### {title}")
    for group in _chunked(fields, columns):
        cols = st.columns(len(group))
        for widget, field in zip(cols, group):
            attr = field["attr"]
            label = field["label"]
            dtype = field.get("type", "float")
            min_value = field.get("min")
            max_value = field.get("max")
            step = field.get("step")
            fmt = field.get("format")
            help_text = field.get("help")

            default_val = getattr(defaults, attr)

            if min_value is not None:
                default_val = max(default_val, min_value)
            if max_value is not None:
                default_val = min(default_val, max_value)
            input_kwargs: Dict[str, Any] = {}
            if help_text:
                input_kwargs["help"] = help_text

            if dtype == "int":
                input_kwargs["value"] = int(round(default_val))
                input_kwargs["step"] = int(step) if step is not None else 1
                if min_value is not None:
                    input_kwargs["min_value"] = int(min_value)
                if max_value is not None:
                    input_kwargs["max_value"] = int(max_value)
                values[attr] = widget.number_input(label, **input_kwargs)
            else:
                input_kwargs["value"] = float(default_val)
                if min_value is not None:
                    input_kwargs["min_value"] = float(min_value)
                if max_value is not None:
                    input_kwargs["max_value"] = float(max_value)
                if step is not None:
                    input_kwargs["step"] = float(step)
                if fmt is not None:
                    input_kwargs["format"] = fmt
                values[attr] = widget.number_input(label, **input_kwargs)


def _render_schedule_editor(
    title: str,
    schedule_key: str,
    default_rows: List[Dict[str, Any]],
    scenario: str,
    namespace: str,
    *,
    fixed_columns: Optional[Dict[str, Any]] = None,
    row_defaults: Optional[Dict[str, Any]] = None,
    allow_yearly_increment: bool = True,
    auto_update_revenue: bool = False,
) -> pd.DataFrame:
    """Render an editable schedule with add/remove and yearly increment controls.

    Parameters
    ----------
    title, schedule_key
        Identifiers for the UI caption and the per-schedule state slot.
    default_rows
        List of default row dictionaries (or dataclasses) that seed the editor
        when no user overrides have been captured for the active scenario.
    scenario, namespace
        Keys used to scope the stored state inside ``st.session_state``.
    fixed_columns
        Optional mapping of column names to constant values. These columns are
        enforced after every edit and rendered as read-only in the table.
    row_defaults
        Optional template applied when users click “Add row”.
    allow_yearly_increment
        Enables the yearly growth controls when the schedule contains numeric
        columns.
    auto_update_revenue
        Recomputes ``Revenue`` from ``Units`` × ``Unit price`` after edits.
    """

    st.markdown(f"**{title}**")
    df, state = _initialise_schedule_state(namespace, scenario, schedule_key, default_rows)
    df = df.convert_dtypes()
    original_columns = [
        col
        for col in df.columns
        if col not in {ROW_REMOVAL_COLUMN, ROW_EDIT_COLUMN}
    ]

    fixed_columns = fixed_columns or {}
    row_defaults = row_defaults or {}

    edit_toggle_key = f"edit_enabled_{namespace}_{schedule_key}_{scenario}"
    edit_enabled = st.checkbox(
        "Enable editing",
        value=st.session_state.get(edit_toggle_key, True),
        key=edit_toggle_key,
        help="Toggle to edit this schedule. When disabled the schedule is read-only.",
    )

    operation_applied = False
    if edit_enabled:
        if ROW_REMOVAL_COLUMN not in df.columns:
            df[ROW_REMOVAL_COLUMN] = False
        if ROW_EDIT_COLUMN not in df.columns:
            df[ROW_EDIT_COLUMN] = False
        controls = st.columns(2)
        if controls[0].button(
            "Add row",
            key=f"add_{namespace}_{schedule_key}_{scenario}",
        ):
            df = _add_schedule_row(df, row_defaults=row_defaults, fixed_columns=fixed_columns)
            operation_applied = True
        if controls[1].button(
            "Remove row",
            key=f"remove_{namespace}_{schedule_key}_{scenario}",
        ):
            if not df.empty:
                df = df.iloc[:-1].reset_index(drop=True)
                operation_applied = True

        df = _ensure_fixed_columns(df, fixed_columns)

        numeric_columns = [
            column
            for column in df.columns
            if pd.to_numeric(df[column], errors="coerce").notna().sum() > 0
        ]

        if allow_yearly_increment and numeric_columns:
            inc_cols, inc_rate_col, inc_btn_col = st.columns([2, 1, 1])
            selected_columns = inc_cols.multiselect(
                "Yearly increment columns",
                options=numeric_columns,
                default=numeric_columns,
                key=f"inc_cols_{namespace}_{schedule_key}_{scenario}",
                help="Select numeric columns that should follow the yearly increment growth.",
            )
            increment_rate = inc_rate_col.number_input(
                "Yearly increment (%)",
                min_value=-100.0,
                max_value=100.0,
                value=0.0,
                step=0.5,
                key=f"inc_rate_{namespace}_{schedule_key}_{scenario}",
            )
            if inc_btn_col.button(
                "Apply yearly increment",
                key=f"apply_inc_{namespace}_{schedule_key}_{scenario}",
            ):
                if selected_columns:
                    df = _apply_yearly_increment(df, selected_columns, increment_rate / 100.0)
                    operation_applied = True
                else:
                    st.warning("Select at least one column before applying an increment.")
        elif allow_yearly_increment:
            st.caption("No numeric columns available for yearly increment.")
    else:
        if ROW_REMOVAL_COLUMN in df.columns:
            df = df.drop(columns=[ROW_REMOVAL_COLUMN])
        if ROW_EDIT_COLUMN in df.columns:
            df = df.drop(columns=[ROW_EDIT_COLUMN])
        st.caption("Enable editing to modify this schedule. Current values remain read-only.")

    if (operation_applied and auto_update_revenue) or (auto_update_revenue and edit_enabled):
        df = _auto_compute_revenue(df)

    df = _ensure_fixed_columns(df, fixed_columns)

    column_config = _build_column_config(
        df,
        disabled=not edit_enabled,
        fixed=fixed_columns,
    )
    instructions: List[str] = []
    if edit_enabled and ROW_EDIT_COLUMN in df.columns:
        column_config[ROW_EDIT_COLUMN] = st.column_config.CheckboxColumn(
            "Edit",
            help="Tick to open the inline row editor before saving.",
            disabled=False,
        )
        instructions.append("Tick 'Edit' to open the row editor and save updates below.")
    if edit_enabled and ROW_REMOVAL_COLUMN in df.columns:
        column_config[ROW_REMOVAL_COLUMN] = st.column_config.CheckboxColumn(
            "Remove",
            help="Tick to remove this row when saving edits.",
            disabled=False,
        )
        instructions.append("Use the 'Remove' checkbox to delete the row before saving.")
    if instructions:
        st.caption(" ".join(instructions))

    try:
        edited_df = st.data_editor(
            df,
            num_rows="dynamic",
            use_container_width=True,
            hide_index=True,
            column_config=column_config,
            disabled=not edit_enabled,
            key=f"editor_{namespace}_{schedule_key}_{scenario}",
        )
    except StreamlitAPIException:
        fallback_df = df.copy()
        for col in fallback_df.columns:
            if col in {ROW_EDIT_COLUMN, ROW_REMOVAL_COLUMN}:
                continue
            if col in fixed_columns:
                continue
            series = fallback_df[col]
            if series.dtype == object:
                observed = {type(v) for v in series.dropna().tolist()}
                if len(observed) > 1:
                    fallback_df[col] = series.apply(
                        lambda value: "" if pd.isna(value) else str(value)
                    )
        st.warning(
            "Some columns had mixed data types. Falling back to text-safe editing "
            "for compatibility."
        )
        edited_df = st.data_editor(
            fallback_df,
            num_rows="dynamic",
            use_container_width=True,
            hide_index=True,
            disabled=not edit_enabled,
            key=f"editor_fallback_{namespace}_{schedule_key}_{scenario}",
        )

    edited_df = pd.DataFrame(edited_df)
    row_form_applied = False
    if ROW_EDIT_COLUMN in edited_df.columns:
        edit_mask = edited_df[ROW_EDIT_COLUMN].fillna(False)
        edit_indices = [idx for idx, flag in enumerate(edit_mask) if flag]
        for row_idx in edit_indices:
            edited_df, submitted = _render_row_edit_form(
                edited_df,
                row_idx,
                fixed_columns,
                namespace,
                schedule_key,
                scenario,
            )
            row_form_applied = row_form_applied or submitted
        edited_df[ROW_EDIT_COLUMN] = False
    if ROW_REMOVAL_COLUMN in edited_df.columns:
        remove_mask = edited_df[ROW_REMOVAL_COLUMN].fillna(False)
        if remove_mask.any():
            edited_df = edited_df.loc[~remove_mask].reset_index(drop=True)
        edited_df = edited_df.drop(columns=[ROW_REMOVAL_COLUMN])
    if ROW_EDIT_COLUMN in edited_df.columns:
        edited_df = edited_df.drop(columns=[ROW_EDIT_COLUMN])
    if row_form_applied:
        operation_applied = True
    if auto_update_revenue:
        edited_df = _auto_compute_revenue(edited_df)
    edited_df = _ensure_fixed_columns(edited_df, fixed_columns)

    state["data"] = edited_df.replace({pd.NA: None}).to_dict("records")
    return edited_df.reindex(columns=original_columns, fill_value=None)


def _render_analytics_schedule(
    title: str,
    schedule_id: str,
    df: pd.DataFrame,
    scenario: str,
    *,
    namespace: str = "advanced_schedule_state",
    allow_yearly_increment: bool = False,
    row_defaults: Optional[Dict[str, Any]] = None,
    fixed_columns: Optional[Dict[str, Any]] = None,
    iterable_columns: Optional[List[str]] = None,
) -> pd.DataFrame:
    """Convenience wrapper to expose schedule editing in advanced analytics."""

    if df is None or df.empty:
        return df

    clean_df = df.copy()
    iterable_columns = iterable_columns or []
    if iterable_columns:
        clean_df = _stringify_iterable_columns(clean_df, iterable_columns)

    defaults = clean_df.where(pd.notna(clean_df), None).to_dict("records")
    template = row_defaults or {column: None for column in clean_df.columns}
    if iterable_columns and template:
        template = template.copy()
        for column in iterable_columns:
            if column in template:
                template[column] = _stringify_iterable_value(template[column])

    return _render_schedule_editor(
        title,
        _sanitize_key(schedule_id or title),
        defaults,
        scenario,
        namespace=namespace,
        fixed_columns=fixed_columns,
        row_defaults=template,
        allow_yearly_increment=allow_yearly_increment,
    ).pipe(lambda edited: _parse_iterable_columns(edited, iterable_columns) if iterable_columns else edited)


def _render_ai_settings(payload: dict, container: Optional[DeltaGenerator] = None) -> None:
    """Render AI and machine-learning configuration controls."""

    target = container or st
    settings = st.session_state.setdefault("ai_settings", _payload_to_ai_settings(payload))
    st.session_state.setdefault("ai_api_key", settings.get("api_key", ""))

    provider_options = list(AI_PROVIDER_OPTIONS)
    if settings.get("provider") not in provider_options:
        provider_options.append(settings.get("provider"))

    current_provider = settings.get("provider", "OpenAI")
    try:
        provider_index = provider_options.index(current_provider)
    except ValueError:
        provider_index = 0

    ml_defaults = [
        ML_METHOD_LABELS.get(code, code.replace("_", " ").title())
        for code in settings.get("ml_methods", ["linear_regression"])
    ]
    feature_defaults = [
        GEN_AI_FEATURE_LABELS.get(code, code.replace("_", " ").title())
        for code in settings.get("generative_features", ["summary"])
    ]

    form = target.form("ai_settings_form")
    with form:
        enabled = form.checkbox(
            "Enable AI Enhancements",
            value=bool(settings.get("enabled", False)),
            help="Toggle machine-learning forecasts and generative commentary.",
        )
        provider = form.selectbox(
            "Provider",
            provider_options,
            index=provider_index,
            help="Select the API provider powering generative insights.",
        )
        model = form.text_input(
            "Model",
            value=settings.get("model", "gpt-4"),
            help="Name of the deployed model (for example `gpt-4o-mini`).",
        )
        horizon = form.number_input(
            "Forecast Horizon (years)",
            min_value=0,
            max_value=20,
            value=int(settings.get("forecast_horizon", 3)),
            step=1,
            help="Number of additional years used for machine-learning revenue forecasts.",
        )
        auto_run_predictive = form.checkbox(
            "Auto-run predictive analytics",
            value=bool(settings.get("auto_run_predictive", False)),
            help="Automatically refresh predictive analytics when inputs are updated.",
        )

        ml_selection = form.multiselect(
            "Machine Learning Methods",
            list(ML_METHOD_LABELS.values()),
            default=ml_defaults,
            help="Choose algorithms applied to projected net revenue.",
        )
        feature_selection = form.multiselect(
            "Generative Features",
            list(GEN_AI_FEATURE_LABELS.values()),
            default=feature_defaults,
            help="Pick the narrative focus areas generated by the AI summary.",
        )
        api_key = form.text_input(
            "API Key",
            value=st.session_state.get("ai_api_key", ""),
            type="password",
            help="Store your provider API key securely. Keys are retained only for the current session.",
        )

        submitted = form.form_submit_button("Save AI Configuration")

    if submitted:
        ml_codes = [
            ML_LABEL_TO_CODE.get(label, label.replace(" ", "_").lower()) for label in ml_selection
        ]
        feature_codes = [
            GEN_AI_LABEL_TO_CODE.get(label, label.replace(" ", "_").lower())
            for label in feature_selection
        ]

        settings.update(
            {
                "enabled": enabled,
                "provider": provider,
                "model": model.strip() or "gpt-4",
                "forecast_horizon": int(horizon),
                "auto_run_predictive": auto_run_predictive,
                "ml_methods": ml_codes or ["linear_regression"],
                "generative_features": feature_codes or ["summary"],
                "api_key": api_key.strip(),
            }
        )
        st.session_state["ai_settings"] = settings
        st.session_state["ai_api_key"] = settings.get("api_key", "")
        _ai_settings_to_payload(settings, payload)
        st.success("AI configuration updated. Rerunning the model with the new settings.")
        _rerun()


def assumptions_form(defaults: Assumptions) -> Assumptions:
    st.header("Model assumptions")

    farm_name = st.text_input("Farm name", defaults.farm_name)

    st.subheader("Input Landing Page")
    st.caption("Adjust the production, pricing, cost, and capital structure assumptions below.")

    values: Dict[str, Any] = {}

    timeline_placeholder = st.container()

    _render_input_section(
        "Production",
        [
            {"attr": "cycles_per_year", "label": "Cycles per year", "min": 1, "max": 12, "type": "int"},
            {
                "attr": "production_start_year",
                "label": "Production start year",
                "min": 1900,
                "max": 2100,
                "type": "int",
            },
            {
                "attr": "production_horizon_years",
                "label": "Production horizon (years)",
                "min": 1,
                "max": 40,
                "type": "int",
            },
            {
                "attr": "birds_per_cycle",
                "label": "Birds per cycle",
                "min": 1000,
                "max": 100000,
                "step": 1000,
                "type": "int",
            },
            {
                "attr": "mortality_rate",
                "label": "Mortality rate",
                "min": 0.0,
                "max": 0.2,
                "step": 0.005,
                "format": "%.3f",
            },
            {
                "attr": "final_weight_kg",
                "label": "Final weight (kg)",
                "min": 1.0,
                "max": 4.0,
                "step": 0.1,
            },
        ],
        defaults,
        values,
        columns=4,
    )

    _render_input_section(
        "Pricing",
        [
            {
                "attr": "live_price_per_kg",
                "label": "Live broiler price per kg",
                "min": 0.5,
                "max": 5.0,
                "step": 0.05,
            },
            {
                "attr": "eggs_price_per_dozen",
                "label": "Eggs price per dozen",
                "min": 0.5,
                "max": 8.0,
                "step": 0.1,
            },
            {
                "attr": "eggs_per_bird_per_cycle",
                "label": "Eggs per bird per cycle",
                "min": 0.0,
                "max": 1000.0,
                "step": 1.0,
            },
            {
                "attr": "manure_price_per_ton",
                "label": "Poultry manure price per ton",
                "min": 0.0,
                "max": 150.0,
                "step": 1.0,
            },
            {
                "attr": "live_bird_price_per_head",
                "label": "Live bird price per head",
                "min": 0.0,
                "max": 10.0,
                "step": 0.1,
            },
            {
                "attr": "byproduct_price_per_kg",
                "label": "By-product price per kg",
                "min": 0.0,
                "max": 5.0,
                "step": 0.05,
            },
            {
                "attr": "price_growth",
                "label": "Annual price growth",
                "min": -0.05,
                "max": 0.1,
                "step": 0.005,
                "format": "%.3f",
            },
        ],
        defaults,
        values,
        columns=3,
    )

    _render_input_section(
        "Operating costs",
        [
            {
                "attr": "feed_conversion_ratio",
                "label": "Feed conversion ratio",
                "min": 1.0,
                "max": 2.5,
                "step": 0.05,
            },
            {
                "attr": "feed_cost_per_kg",
                "label": "Feed cost per kg",
                "min": 0.2,
                "max": 1.0,
                "step": 0.01,
            },
            {
                "attr": "chick_cost",
                "label": "Chick cost",
                "min": 0.3,
                "max": 2.0,
                "step": 0.05,
            },
            {
                "attr": "processing_cost_per_bird",
                "label": "Processing cost per bird",
                "min": 0.05,
                "max": 1.0,
                "step": 0.05,
            },
            {
                "attr": "vaccination_cost_per_bird",
                "label": "Vaccination cost per bird",
                "min": 0.0,
                "max": 0.5,
                "step": 0.01,
            },
            {
                "attr": "litter_disposal_per_cycle",
                "label": "Litter & disposal per cycle",
                "min": 0.0,
                "max": 10000.0,
                "step": 100.0,
            },
            {
                "attr": "propane_per_cycle",
                "label": "Propane per cycle",
                "min": 0.0,
                "max": 20000.0,
                "step": 100.0,
            },
            {
                "attr": "electricity_per_cycle",
                "label": "Electricity per cycle",
                "min": 0.0,
                "max": 10000.0,
                "step": 100.0,
            },
            {
                "attr": "labor_per_cycle",
                "label": "Labor per cycle",
                "min": 0.0,
                "max": 50000.0,
                "step": 500.0,
            },
            {
                "attr": "maintenance_per_cycle",
                "label": "Maintenance per cycle",
                "min": 0.0,
                "max": 10000.0,
                "step": 100.0,
            },
            {
                "attr": "management_fee_per_cycle",
                "label": "Management fee per cycle",
                "min": 0.0,
                "max": 10000.0,
                "step": 100.0,
            },
            {
                "attr": "insurance_per_cycle",
                "label": "Insurance per cycle",
                "min": 0.0,
                "max": 10000.0,
                "step": 100.0,
            },
            {
                "attr": "overhead_per_cycle",
                "label": "Overhead per cycle",
                "min": 0.0,
                "max": 10000.0,
                "step": 100.0,
            },
            {
                "attr": "cost_inflation",
                "label": "Cost inflation",
                "min": -0.05,
                "max": 0.1,
                "step": 0.005,
                "format": "%.3f",
            },
        ],
        defaults,
        values,
        columns=4,
    )

    _render_input_section(
        "Capital & financing",
        [
            {
                "attr": "capex_housing",
                "label": "Housing capex",
                "min": 0.0,
                "max": 5000000.0,
                "step": 10000.0,
            },
            {
                "attr": "capex_equipment",
                "label": "Equipment capex",
                "min": 0.0,
                "max": 2000000.0,
                "step": 5000.0,
            },
            {
                "attr": "ar_days",
                "label": "AR days",
                "min": 0.0,
                "max": 365.0,
                "step": 1.0,
            },
            {
                "attr": "inventory_days",
                "label": "Inventory days",
                "min": 0.0,
                "max": 365.0,
                "step": 1.0,
            },
            {
                "attr": "ap_days",
                "label": "AP days",
                "min": 0.0,
                "max": 365.0,
                "step": 1.0,
            },
            {
                "attr": "depreciation_years",
                "label": "Depreciation years",
                "min": 1,
                "max": 40,
                "type": "int",
            },
            {
                "attr": "maintenance_capex_annual",
                "label": "Maintenance capex (annual)",
                "min": 0.0,
                "max": 200000.0,
                "step": 5000.0,
            },
            {
                "attr": "debt_ratio",
                "label": "Debt ratio",
                "min": 0.0,
                "max": 1.0,
                "step": 0.05,
            },
            {
                "attr": "debt_interest_rate",
                "label": "Debt interest rate",
                "min": 0.0,
                "max": 0.2,
                "step": 0.005,
                "format": "%.3f",
            },
            {
                "attr": "debt_term_years",
                "label": "Debt term (years)",
                "min": 1,
                "max": 30,
                "type": "int",
            },
            {
                "attr": "discount_rate",
                "label": "Discount rate",
                "min": 0.0,
                "max": 0.5,
                "step": 0.01,
                "format": "%.3f",
            },
            {
                "attr": "tax_rate",
                "label": "Tax rate",
                "min": 0.0,
                "max": 0.5,
                "step": 0.01,
                "format": "%.3f",
            },
        ],
        defaults,
        values,
        columns=4,
    )

    assumptions = Assumptions(
        farm_name=farm_name,
        cycles_per_year=int(values["cycles_per_year"]),
        production_start_year=int(values["production_start_year"]),
        production_horizon_years=int(values["production_horizon_years"]),
        birds_per_cycle=int(values["birds_per_cycle"]),
        mortality_rate=float(values["mortality_rate"]),
        final_weight_kg=float(values["final_weight_kg"]),
        live_price_per_kg=float(values["live_price_per_kg"]),
        chick_cost=float(values["chick_cost"]),
        feed_conversion_ratio=float(values["feed_conversion_ratio"]),
        feed_cost_per_kg=float(values["feed_cost_per_kg"]),
        processing_cost_per_bird=float(values["processing_cost_per_bird"]),
        vaccination_cost_per_bird=float(values["vaccination_cost_per_bird"]),
        litter_disposal_per_cycle=float(values["litter_disposal_per_cycle"]),
        propane_per_cycle=float(values["propane_per_cycle"]),
        electricity_per_cycle=float(values["electricity_per_cycle"]),
        labor_per_cycle=float(values["labor_per_cycle"]),
        maintenance_per_cycle=float(values["maintenance_per_cycle"]),
        management_fee_per_cycle=float(values["management_fee_per_cycle"]),
        insurance_per_cycle=float(values["insurance_per_cycle"]),
        overhead_per_cycle=float(values["overhead_per_cycle"]),
        capex_housing=float(values["capex_housing"]),
        capex_equipment=float(values["capex_equipment"]),
        ar_days=float(values["ar_days"]),
        inventory_days=float(values["inventory_days"]),
        ap_days=float(values["ap_days"]),
        working_capital=float(defaults.working_capital),
        discount_rate=float(values["discount_rate"]),
        price_growth=float(values["price_growth"]),
        eggs_price_per_dozen=float(values["eggs_price_per_dozen"]),
        eggs_per_bird_per_cycle=float(values["eggs_per_bird_per_cycle"]),
        manure_price_per_ton=float(values["manure_price_per_ton"]),
        live_bird_price_per_head=float(values["live_bird_price_per_head"]),
        byproduct_price_per_kg=float(values["byproduct_price_per_kg"]),
        cost_inflation=float(values["cost_inflation"]),
        tax_rate=float(values["tax_rate"]),
        debt_ratio=float(values["debt_ratio"]),
        debt_interest_rate=float(values["debt_interest_rate"]),
        debt_term_years=int(values["debt_term_years"]),
        depreciation_years=int(values["depreciation_years"]),
        maintenance_capex_annual=float(values["maintenance_capex_annual"]),
    )

    with timeline_placeholder:
        start_year = int(assumptions.production_start_year)
        horizon_years = int(assumptions.production_horizon_years)
        if horizon_years <= 0:
            horizon_years = 1
        end_year = start_year + horizon_years - 1
        timeline_cols = st.columns(2)
        timeline_cols[0].metric("Production start year", start_year)
        timeline_cols[1].metric("Production end year", end_year)

    return assumptions
def main() -> None:
    st.set_page_config(page_title="Broiler Chicken Financial Model", layout="wide")
    st.title("Broiler Chicken Financial Model")
    st.caption("Interactively explore production, operating, and financing assumptions for a broiler chicken farm.")

    st.subheader("Scenario selection")
    selected_scenario = st.selectbox("Scenario", SCENARIO_OPTIONS, key="selected_scenario")

    scenario_store = st.session_state.setdefault("scenario_store", {})
    if selected_scenario not in scenario_store:
        scenario_store[selected_scenario] = {
            "assumptions": asdict(Assumptions()),
            "ai_settings": DEFAULT_AI_SETTINGS.copy(),
            "investor_benchmarks": DEFAULT_INVESTOR_BENCHMARKS.copy(),
        }
    payload = _normalize_scenario_payload(scenario_store.get(selected_scenario))
    scenario_store[selected_scenario] = payload

    defaults = Assumptions(**payload["assumptions"])

    (
        input_tab,
        production_tab,
        financials_tab,
        analytics_tab,
        ai_ml_tab,
    ) = st.tabs(
        [
            "Input Landing Page",
            "Production & revenues",
            "Financial statements",
            "Advanced analytics",
            "AI & machine learning settings",
        ]
    )

    with input_tab:
        assumptions = assumptions_form(defaults)

    previous_assumptions = copy.deepcopy(payload.get("assumptions", {}))
    payload["assumptions"] = asdict(assumptions)
    if previous_assumptions != payload["assumptions"]:
        _reset_scenario_edit_states(selected_scenario)
        governance_store = st.session_state.setdefault("governance_log", {})
        scenario_log = governance_store.setdefault(selected_scenario, [])
        scenario_log.append(
            {
                "timestamp": datetime.now(timezone.utc).isoformat(),
                "event": "Assumptions updated",
            }
        )
        st.session_state.governance_log = governance_store

    ai_settings = _payload_to_ai_settings(payload)
    payload["ai_settings"] = ai_settings
    st.session_state["ai_settings"] = ai_settings

    input_page = copy.deepcopy(payload)
    snapshot = st.session_state.get("input_snapshot")
    if snapshot is None or st.session_state.get("snapshot_scenario") != selected_scenario:
        snapshot = copy.deepcopy(input_page)
        st.session_state.snapshot_scenario = selected_scenario
    elif snapshot != input_page:
        snapshot = copy.deepcopy(input_page)
    st.session_state.input_snapshot = snapshot

    model, results = _ensure_scenario_payload(selected_scenario, snapshot)
    model.scenario = selected_scenario
    st.session_state.model_results = (model, results)

    valuation = results["valuation"]
    assumption_schedule_df = pd.DataFrame(results["assumptions_schedule"])
    revenue_schedules = results["revenue_schedules"]
    revenue_summary = results.get(
        "revenue_summary",
        summarise_revenue_totals(
            revenue_schedules,
            model.assumptions.cycles_per_year,
            model.assumptions.production_horizon_years,
            model.assumptions.production_start_year,
        ),
    )
    financials = results["financial_statements"]
    advanced = results["advanced_analytics"]

    col1, col2, col3 = st.columns(3)
    col1.metric("NPV", f"${valuation['npv']:,.0f}")
    col2.metric("IRR", f"{valuation['irr']:.2%}")
    col3.metric("Discount rate", f"{valuation['discount_rate']:.2%}")

    cycles_df = pd.DataFrame([asdict(cycle) for cycle in results["cycles"]])
    annual_summary_obj = results["annual"]
    annual_df = pd.DataFrame([asdict(annual_summary_obj)])
    valuation_df = pd.DataFrame([valuation])
    cashflow_df = pd.DataFrame([asdict(row) for row in results["cashflows"]])
    cashflow_bridge_frames = _prepare_cashflow_bridge_frames(cashflow_df)
    income_df = pd.DataFrame([asdict(row) for row in financials["income_statement"]])
    balance_df = pd.DataFrame([asdict(row) for row in financials["balance_sheet"]])
    cash_statement_df = pd.DataFrame([asdict(row) for row in financials["cash_flow_statement"]])
    loan_df = pd.DataFrame(financials["loan_schedule"])
    asset_schedule_df = (
        balance_df[
            ["year", "cash", "working_capital", "net_ppe", "total_assets"]
        ]
        .copy()
        .rename(
            columns={
                "year": "Year",
                "cash": "Ending cash",
                "working_capital": "Working capital",
                "net_ppe": "Net PP&E",
                "total_assets": "Total assets",
            }
        )
    )
    metrics_df = pd.DataFrame(advanced["metrics"])
    dscr_df = pd.DataFrame(advanced["dscr"])
    trend_df = pd.DataFrame(advanced["trend"])
    returns_df = pd.DataFrame(advanced.get("returns", []))
    coverage_df = pd.DataFrame(advanced.get("coverage", []))
    leverage_df = pd.DataFrame(advanced.get("leverage", []))
    what_if_df = pd.DataFrame(advanced.get("what_if", []))
    monte_carlo = advanced.get("monte_carlo", {})
    monte_carlo_summary_df = pd.DataFrame([monte_carlo["summary"]]) if monte_carlo.get("summary") else pd.DataFrame()
    monte_carlo_samples_df = pd.DataFrame(monte_carlo.get("samples", []))
    break_even_df = pd.DataFrame(advanced.get("break_even", []))
    goal_seek = advanced.get("goal_seek", {})
    predictive = advanced.get("predictive", {})
    forecast_df = pd.DataFrame(predictive.get("automated_forecast", []))
    time_series_df = pd.DataFrame(predictive.get("time_series", {}).get("forecast", []))
    risk_metadata = predictive.get("risk_anomalies", {})
    risk_df = pd.DataFrame(risk_metadata.get("observations", []))
    ml_methods_df = pd.DataFrame(predictive.get("ml_methods", []))
    scenario_df = pd.DataFrame(advanced.get("scenario_planning", []))
    custom_simulations_data = advanced.get("custom_simulations", {})
    custom_definition_defaults = advanced.get("custom_simulation_definitions")
    if not custom_definition_defaults:
        custom_definition_defaults = copy.deepcopy(
            DEFAULT_CUSTOM_SIMULATION_DEFINITIONS
        )
    else:
        custom_definition_defaults = copy.deepcopy(custom_definition_defaults)

    metrics_lookup = metrics_df.set_index("metric")["value"] if not metrics_df.empty else pd.Series(dtype="float64")

    def _metric_value(name: str) -> float:
        try:
            return float(metrics_lookup.get(name, float("nan")))
        except Exception:
            return float("nan")

    base_metrics = advanced.get("base_metrics") or {
        "npv": valuation.get("npv"),
        "irr": valuation.get("irr"),
        "avg_dscr": _metric_value("Average DSCR"),
        "min_dscr": _metric_value("Minimum DSCR"),
        "payback": _metric_value("Payback period (years)"),
    }

    payback_period_value = _to_float(_metric_value("Payback period (years)"))
    if isinstance(payback_period_value, float) and math.isnan(payback_period_value):
        payback_period_value = None
    st.markdown("### Cross-page simulation KPIs")
    kpi_cols = st.columns(1)
    kpi_cols[0].metric(
        "Payback period",
        f"{payback_period_value:.2f} years"
        if payback_period_value is not None
        else "N/A years",
    )

    with production_tab:
        download_container = st.container()
        export_map: Dict[str, Dict[str, Any]] = st.session_state.setdefault(
            "model_export_map", {}
        )
        export_payload = export_map.get(selected_scenario)

        with download_container:
            st.markdown("### Excel export")
            excel_health = _excel_dependency_health()
            health_badge = (
                "🟢 Healthy"
                if excel_health.get("ready")
                else "🟡 Fallback mode (Excel-compatible XML)"
            )
            st.caption(
                "Excel dependency health check: "
                f"{health_badge} | "
                f"xlsxwriter={'✅' if excel_health.get('xlsxwriter') else '❌'} "
                f"openpyxl={'✅' if excel_health.get('openpyxl') else '❌'}"
            )
            if not export_payload:
                if st.button(
                    "Prepare Excel Model",
                    key=f"prepare_excel_{selected_scenario.lower()}",
                ):
                    try:
                        with st.spinner("Preparing Excel workbook..."):
                            excel_bytes = _generate_excel_bytes(
                                model, results, selected_scenario
                            )
                    except RuntimeError as exc:
                        if "Missing Excel writer dependency" in str(exc):
                            with st.spinner("Preparing Excel-compatible fallback export..."):
                                excel_xml_bytes = _generate_excel_xml_bytes(
                                    model, results, selected_scenario
                                )
                            export_payload = {
                                "data": excel_xml_bytes,
                                "file_name": f"Broiler_Financial_Model_{selected_scenario.replace(' ', '_')}_Compatible.xml",
                                "mime": "application/xml",
                                "label": "Download Excel-Compatible Workbook",
                                "fallback": True,
                            }
                        else:
                            export_payload = None
                    else:
                        export_payload = {
                            "data": excel_bytes,
                            "file_name": f"Broiler_Financial_Model_{selected_scenario.replace(' ', '_')}.xlsx",
                            "mime": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                            "label": "Download Excel Model",
                            "fallback": False,
                        }

                    if export_payload:
                        export_map[selected_scenario] = export_payload
                        st.session_state.model_export_map = export_map
            if export_payload:
                if export_payload.get("fallback"):
                    st.warning(
                        "Excel writer dependencies are unavailable. Downloading an Excel-compatible XML workbook instead."
                    )
                st.download_button(
                    export_payload.get("label", "Download Model Export"),
                    data=export_payload["data"],
                    file_name=export_payload["file_name"],
                    mime=export_payload["mime"],
                    key=f"download_excel_{selected_scenario.lower()}",
                )
                if st.button(
                    "Clear Prepared Excel",
                    key=f"clear_excel_{selected_scenario.lower()}",
                ):
                    export_map.pop(selected_scenario, None)
                    st.session_state.model_export_map = export_map
                    export_payload = None
            if not export_payload:
                st.info("Click 'Prepare Excel Model' to generate the workbook for download.")

        st.subheader("Assumptions summary")
        updated_assumptions: List[Dict[str, Any]] = []
        for schedule_name, group in assumption_schedule_df.groupby("schedule", sort=False):
            schedule_key = _sanitize_key(schedule_name)
            defaults = group.drop(columns=["schedule"]).to_dict("records")
            edited_df = _render_schedule_editor(
                schedule_name,
                schedule_key,
                defaults,
                selected_scenario,
                namespace="assumption_schedule_state",
                allow_yearly_increment=True,
            )
            for record in edited_df.replace({pd.NA: None}).to_dict("records"):
                merged = {"schedule": schedule_name}
                merged.update(record)
                updated_assumptions.append(merged)
        if updated_assumptions:
            results["assumptions_schedule"] = updated_assumptions
            new_assumptions = _assumptions_from_schedule(
                updated_assumptions, model.assumptions
            )
            if new_assumptions is not None:
                prior_assumptions = copy.deepcopy(payload.get("assumptions", {}))
                payload["assumptions"] = asdict(new_assumptions)
                scenario_store[selected_scenario]["assumptions"] = asdict(
                    new_assumptions
                )
                if prior_assumptions != payload["assumptions"]:
                    governance_store = st.session_state.setdefault("governance_log", {})
                    scenario_log = governance_store.setdefault(selected_scenario, [])
                    scenario_log.append(
                        {
                            "timestamp": datetime.now(timezone.utc).isoformat(),
                            "event": "Assumption schedule edits applied",
                        }
                    )
                    st.session_state.governance_log = governance_store
                st.session_state.snapshot_scenario = selected_scenario
                st.session_state.input_snapshot = copy.deepcopy(payload)

                updated_model = ScenarioModel(
                    scenario=selected_scenario, assumptions=new_assumptions
                )
                updated_results = generate_model_outputs(
                    new_assumptions,
                    analytics_plan=AnalyticsPlan.summary(),
                )

                scenario_payloads = st.session_state.setdefault(
                    "scenario_payloads", {}
                )
                scenario_payloads[selected_scenario] = {
                    "snapshot": copy.deepcopy(st.session_state.input_snapshot),
                    "model": updated_model,
                    "results": updated_results,
                }
                results = updated_results
                st.session_state.model_results = (updated_model, updated_results)
                st.success(
                    "Assumption schedule edits applied. Recalculating model outputs..."
                )
                _rerun()

        st.subheader("Revenue schedules")
        updated_revenue: Dict[str, List[Dict[str, Any]]] = {}
        for category, rows in revenue_schedules.items():
            schedule_key = _sanitize_key(category)
            defaults = copy.deepcopy(rows)
            row_defaults: Dict[str, Any] = {}
            if rows:
                row_defaults["Notes"] = rows[0].get("Notes")
            edited_df = _render_schedule_editor(
                category,
                schedule_key,
                defaults,
                selected_scenario,
                namespace="revenue_schedule_state",
                fixed_columns={"Category": category},
                row_defaults=row_defaults,
                allow_yearly_increment=True,
                auto_update_revenue=True,
            )
            updated_revenue[category] = (
                edited_df.replace({pd.NA: None}).to_dict("records")
            )
        if updated_revenue:
            revenue_schedules = updated_revenue
            results["revenue_schedules"] = updated_revenue
            revenue_summary = summarise_revenue_totals(
                revenue_schedules,
                model.assumptions.cycles_per_year,
                model.assumptions.production_horizon_years,
                model.assumptions.production_start_year,
            )
            results["revenue_summary"] = revenue_summary

        summary_by_category = pd.DataFrame(revenue_summary.get("by_category", []))
        if not summary_by_category.empty:
            st.markdown("##### Annual revenue by category")
            st.dataframe(
                summary_by_category,
                use_container_width=True,
                hide_index=True,
            )
            revenue_chart = _build_revenue_stack_chart(summary_by_category)
            if revenue_chart is not None:
                st.altair_chart(revenue_chart, use_container_width=True)

        annual_totals_df = pd.DataFrame(revenue_summary.get("annual_totals", []))
        if not annual_totals_df.empty:
            st.markdown("##### Total revenue by year")
            st.dataframe(
                annual_totals_df,
                use_container_width=True,
                hide_index=True,
            )

        st.subheader("Production cycles")
        st.dataframe(cycles_df, use_container_width=True)

        st.subheader("Annual summary")
        st.dataframe(annual_df, use_container_width=True)

        st.subheader("Discounted cash flows")
        st.dataframe(cashflow_df, use_container_width=True)
        if cashflow_bridge_frames:
            waterfall_years = sorted(cashflow_bridge_frames)
            waterfall_year = st.selectbox(
                "Cash-flow bridge year",
                waterfall_years,
                key=f"waterfall_year_{selected_scenario}",
            )
            waterfall_chart = _cashflow_waterfall_chart(
                cashflow_bridge_frames.get(waterfall_year, pd.DataFrame())
            )
            if waterfall_chart is not None:
                st.altair_chart(waterfall_chart, use_container_width=True)

        st.subheader("Debt schedule")
        st.dataframe(loan_df, use_container_width=True, hide_index=True)

        st.subheader("Asset schedule")
        st.dataframe(asset_schedule_df, use_container_width=True, hide_index=True)

    with financials_tab:
        fin_tab1, fin_tab2, fin_tab3, fin_tab4 = st.tabs(
            [
                "Income statement",
                "Balance sheet",
                "Cash flow statement",
                "Debt schedule",
            ]
        )
        with fin_tab1:
            st.dataframe(income_df, use_container_width=True, hide_index=True)
        with fin_tab2:
            st.dataframe(balance_df, use_container_width=True, hide_index=True)
        with fin_tab3:
            st.dataframe(cash_statement_df, use_container_width=True, hide_index=True)
        with fin_tab4:
            st.dataframe(loan_df, use_container_width=True, hide_index=True)


    with ai_ml_tab:
        st.header("Unified AI System")
        st.caption(
            "Single integrated AI workflow combining settings, investor benchmarking, governance, "
            "RAG research, model Q&A, and business-plan generation."
        )
        _render_ai_settings(payload)
        st.markdown("### Unified configuration")
        benchmark_defaults = payload.get(
            "investor_benchmarks", DEFAULT_INVESTOR_BENCHMARKS.copy()
        )
        unified_ai_state: Dict[str, Any] = {
            "config": {
                "scenario": selected_scenario,
                "ai_settings": payload.get("ai_settings", DEFAULT_AI_SETTINGS.copy()),
            },
            "inputs": {},
            "outputs": {},
        }
        bcol1, bcol2, bcol3 = st.columns(3)
        target_irr = bcol1.number_input(
            "Target IRR",
            min_value=0.0,
            max_value=1.0,
            value=float(benchmark_defaults.get("target_irr", 0.18)),
            step=0.01,
            format="%.2f",
        )
        min_dscr = bcol2.number_input(
            "Minimum DSCR",
            min_value=0.5,
            max_value=5.0,
            value=float(benchmark_defaults.get("min_dscr", 1.5)),
            step=0.1,
            format="%.2f",
        )
        max_payback = bcol3.number_input(
            "Max payback (years)",
            min_value=1.0,
            max_value=20.0,
            value=float(benchmark_defaults.get("max_payback_years", 6.0)),
            step=0.5,
            format="%.1f",
        )
        payload["investor_benchmarks"] = {
            "target_irr": float(target_irr),
            "min_dscr": float(min_dscr),
            "max_payback_years": float(max_payback),
        }
        unified_ai_state["config"]["investor_benchmarks"] = payload["investor_benchmarks"]
        st.session_state["scenario_store"][selected_scenario] = payload

        scorecard_df = _build_investor_scorecard(
            valuation,
            dscr_df,
            break_even_df,
            assumptions,
            payload["investor_benchmarks"],
        )
        benchmark_map = {
            "IRR": target_irr,
            "Average DSCR": min_dscr,
            "Payback (years)": max_payback,
        }
        if not scorecard_df.empty:
            scorecard_df = scorecard_df.copy()
            scorecard_df["Benchmark"] = scorecard_df["Metric"].map(benchmark_map)
            scorecard_df["Benchmark"] = scorecard_df["Benchmark"].fillna("N/A")
            st.markdown("### Unified output: investor scorecard")
            st.dataframe(scorecard_df, use_container_width=True, hide_index=True)
        unified_ai_state["outputs"]["investor_scorecard_rows"] = int(
            len(scorecard_df.index) if isinstance(scorecard_df, pd.DataFrame) else 0
        )

        st.markdown("### Unified output: assumption confidence bands")
        confidence_df = _build_assumption_confidence_frame(assumptions)
        st.dataframe(confidence_df, use_container_width=True, hide_index=True)
        unified_ai_state["outputs"]["confidence_band_rows"] = int(
            len(confidence_df.index) if isinstance(confidence_df, pd.DataFrame) else 0
        )

        st.markdown("### Unified output: model governance")
        assumptions_hash = hashlib.sha256(
            json.dumps(payload.get("assumptions", {}), sort_keys=True).encode("utf-8")
        ).hexdigest()[:12]
        governance_log = st.session_state.get("governance_log", {}).get(
            selected_scenario, []
        )
        st.write(f"**Model version:** v1.0")
        st.write(f"**Scenario:** {selected_scenario}")
        st.write(f"**Assumptions hash:** `{assumptions_hash}`")
        st.write(
            f"**Last reviewed (UTC):** {datetime.now(timezone.utc).strftime('%Y-%m-%d %H:%M:%S')}"
        )
        if governance_log:
            st.markdown("**Recent changes:**")
            for entry in governance_log[-5:]:
                st.markdown(f"- {entry.get('timestamp')}: {entry.get('event')}")
        unified_ai_state["outputs"]["governance"] = {
            "assumptions_hash": assumptions_hash,
            "recent_events": governance_log[-5:],
        }

        st.markdown("### Unified input: RAG research library")
        st.caption(
            "Upload research papers/documents to ground business-plan rewrites with evidence."
        )
        st.caption(
            "Supported indexing formats: PDF, Word, Excel, CSV/TSV, JSON, HTML/XML, PPTX, TXT/MD/LOG, and ZIP bundles."
        )
        uploaded_docs = st.file_uploader(
            "Upload research documents",
            type=_rag_supported_extensions(include_archives=True),
            accept_multiple_files=True,
            key=f"rag_upload_{selected_scenario}",
        )
        rag_store = st.session_state.setdefault("rag_document_store", {})
        rag_store.setdefault(selected_scenario, {"documents": [], "chunks": []})
        if st.button(
            "Index uploaded documents",
            key=f"index_rag_documents_{selected_scenario}",
        ):
            documents: List[Dict[str, Any]] = []
            all_chunks: List[Dict[str, Any]] = []
            parse_diagnostics: List[Dict[str, Any]] = []
            expanded_uploads = _expand_uploaded_documents(uploaded_docs)
            for upload_entry in expanded_uploads:
                file_name = upload_entry["name"]
                file_payload = upload_entry["payload"]
                if not upload_entry.get("supported", True):
                    parse_diagnostics.append(
                        {
                            "File": file_name,
                            "Type": file_name.split(".")[-1].lower() if "." in file_name else "",
                            "Parser": "n/a",
                            "Characters": 0,
                            "Chunks": 0,
                            "Status": "Skipped",
                            "Error": "Unsupported extension for RAG indexing configuration.",
                        }
                    )
                    continue
                extracted, parser_used, parse_error = _extract_text_from_payload(
                    file_name, file_payload
                )
                if not extracted.strip():
                    st.warning(
                        f"Could not parse text for `{file_name}`. "
                        f"Parser: {parser_used}. "
                        f"{parse_error or 'Install optional parser dependencies for this format if needed.'}"
                    )
                    parse_diagnostics.append(
                        {
                            "File": file_name,
                            "Type": file_name.split(".")[-1].lower() if "." in file_name else "",
                            "Parser": parser_used,
                            "Characters": 0,
                            "Chunks": 0,
                            "Status": "Failed",
                            "Error": parse_error or "No text extracted",
                        }
                    )
                    continue
                doc_record = {
                    "name": file_name,
                    "char_count": len(extracted),
                }
                documents.append(doc_record)
                new_chunks = _chunk_document_text(file_name, extracted)
                all_chunks.extend(new_chunks)
                parse_diagnostics.append(
                    {
                        "File": file_name,
                        "Type": file_name.split(".")[-1].lower() if "." in file_name else "",
                        "Parser": parser_used,
                        "Characters": len(extracted),
                        "Chunks": len(new_chunks),
                        "Status": "Indexed",
                        "Error": "",
                    }
                )
            rag_store[selected_scenario] = {
                "documents": documents,
                "chunks": all_chunks,
            }
            st.session_state.rag_document_store = rag_store
            st.success(
                f"Indexed {len(documents)} document(s) into {len(all_chunks)} chunk(s)."
            )
            if parse_diagnostics:
                st.markdown("#### Indexing diagnostics")
                st.dataframe(pd.DataFrame(parse_diagnostics), use_container_width=True, hide_index=True)

        rag_docs = rag_store.get(selected_scenario, {}).get("documents", [])
        rag_chunks = rag_store.get(selected_scenario, {}).get("chunks", [])
        st.write(f"Indexed documents: **{len(rag_docs)}** | Chunks: **{len(rag_chunks)}**")
        unified_ai_state["inputs"]["rag"] = {
            "documents": len(rag_docs),
            "chunks": len(rag_chunks),
        }

        st.markdown("### Unified interaction: model Q&A chatbox")
        st.caption(
            "Ask a question and get an answer strictly from this scenario's model outputs."
        )
        chat_history_key = f"model_chat_history_{selected_scenario}"
        chat_history = st.session_state.setdefault(chat_history_key, [])
        for message in chat_history:
            with st.chat_message(message.get("role", "assistant")):
                st.markdown(message.get("content", ""))

        prompt = st.chat_input(
            "Ask a question about this model scenario...",
            key=f"model_chat_input_{selected_scenario}",
        )
        qa_context_frames: Dict[str, pd.DataFrame] = {
            "Assumptions schedule": assumption_schedule_df,
            "Annual summary": annual_df,
            "Valuation": valuation_df,
            "Cash flows": cashflow_df,
            "Income statement": income_df,
            "Balance sheet": balance_df,
            "Cash flow statement": cash_statement_df,
            "Debt schedule": loan_df,
            "Advanced metrics": metrics_df,
            "DSCR": dscr_df,
            "Trend analysis": trend_df,
            "Coverage analysis": coverage_df,
            "Leverage analysis": leverage_df,
            "What-if analysis": what_if_df,
            "Break-even analysis": break_even_df,
            "Forecast projections": forecast_df,
            "Time-series forecast": time_series_df,
            "Scenario planning": scenario_df,
            "Risk observations": risk_df,
            "ML methods": ml_methods_df,
            "Monte Carlo summary": monte_carlo_summary_df,
            "Monte Carlo samples": monte_carlo_samples_df,
            "Revenue by category": summary_by_category,
            "Revenue annual totals": annual_totals_df,
        }
        for category, rows in revenue_schedules.items():
            qa_context_frames[f"Revenue schedule - {category}"] = pd.DataFrame(rows)
        unified_ai_state["inputs"]["context_tables"] = sorted(qa_context_frames.keys())
        investor_recommendations = _generate_investor_recommendations(
            assumptions,
            valuation,
            annual_df,
            dscr_df,
            break_even_df,
        )
        orchestration_config = {
            "scenario": selected_scenario,
            "model_version": "v1.0",
            "governance_hash": assumptions_hash,
            "payback_years": payback_period_value,
        }
        unified_ai_state["config"]["orchestration"] = orchestration_config
        orchestration_state = _build_unified_orchestration_state(
            orchestration_config,
            qa_context_frames,
            assumptions,
            valuation,
            rag_chunks,
            payload["investor_benchmarks"],
            investor_recommendations,
        )
        readiness = orchestration_state["reasoning"]["investor_readiness"]
        st.markdown("### Unified intelligence engine")
        st.caption(
            "Cross-functional reasoning with shared context, grounded evidence, and investor decision support."
        )
        readiness_cols = st.columns(3)
        readiness_cols[0].metric("Readiness verdict", readiness.get("verdict", "N/A"))
        readiness_cols[1].metric("Readiness score", f"{readiness.get('score', 0.0):.0f}/100")
        readiness_cols[2].metric(
            "Strategic signals",
            len(orchestration_state["reasoning"].get("strategic_analysis", [])),
        )
        with st.expander("Explainable strategic analysis", expanded=False):
            for line in orchestration_state["reasoning"].get("strategic_analysis", []):
                st.markdown(f"- {line}")
            for line in orchestration_state["reasoning"].get(
                "proactive_recommendations", []
            ):
                st.markdown(f"- Recommendation: {line}")
        unified_ai_state["outputs"]["orchestration_readiness"] = readiness
        if prompt:
            chat_history.append({"role": "user", "content": prompt})
            answer = _answer_unified_orchestrator_question(
                prompt, orchestration_state
            )
            chat_history.append({"role": "assistant", "content": answer})
            st.session_state[chat_history_key] = chat_history
            _rerun()

        st.divider()
        st.subheader("Unified output: Business Plan Agent")
        st.caption(
            "Generate a comprehensive, investor-ready business plan with automated analysis "
            "and visual exhibits derived from the live model outputs."
        )
        generate_plan = st.button(
            "Generate full business plan",
            key=f"generate_business_plan_{selected_scenario}",
            type="primary",
        )
        plan_state_key = f"business_plan_ready_{selected_scenario}"
        if generate_plan:
            st.session_state[plan_state_key] = True

        if st.session_state.get(plan_state_key, False):
            base_plan_markdown = _build_business_plan_markdown(
                selected_scenario,
                assumptions,
                valuation,
                annual_df,
                dscr_df,
                break_even_df,
                monte_carlo_summary_df,
            )
            section_queries = {
                "Executive Summary": "investment thesis return profile npv irr",
                "Production & Operating Plan": "broiler operations feed conversion mortality throughput productivity",
                "Financial Performance Plan": "ebitda cash flow dscr payback financing",
                "Market, Pricing, and Break-even Strategy": "market demand price benchmark break-even",
                "Risk Assessment & Mitigation": "risk volatility downside sensitivity scenario",
                "Implementation Roadmap": "timeline milestones implementation plan",
            }
            rewrite_plan = st.button(
                "Rewrite using RAG evidence",
                key=f"rewrite_plan_with_rag_{selected_scenario}",
                disabled=not rag_chunks,
            )
            citation_entries: List[Dict[str, str]] = []
            if rewrite_plan and rag_chunks:
                evidence = _retrieve_evidence_for_sections(rag_chunks, section_queries, top_k=3)
                plan_markdown, citation_entries = _rewrite_plan_with_rag_evidence(
                    base_plan_markdown, evidence
                )
                st.success("Business plan rewritten using indexed RAG evidence.")
            else:
                plan_markdown = base_plan_markdown
            st.markdown(plan_markdown)
            st.markdown("### Investor attractiveness recommendations")
            for recommendation in investor_recommendations:
                st.markdown(f"- {recommendation}")

            ic_pack_text = _build_ic_pack_markdown(
                selected_scenario,
                scorecard_df,
                plan_markdown,
                investor_recommendations,
            )
            st.download_button(
                "Download IC Pack (Markdown)",
                data=ic_pack_text.encode("utf-8"),
                file_name=f"IC_Pack_{selected_scenario.replace(' ', '_')}.md",
                mime="text/markdown",
                key=f"download_ic_pack_{selected_scenario}",
            )
            export_title = f"Business Plan - {selected_scenario}"
            export_cols = st.columns(3)
            with export_cols[0]:
                try:
                    docx_bytes = _generate_business_plan_docx_bytes(
                        export_title,
                        plan_markdown,
                        investor_recommendations,
                        citation_entries,
                    )
                except RuntimeError as exc:
                    st.warning(str(exc))
                else:
                    st.download_button(
                        "Download Word (.docx)",
                        data=docx_bytes,
                        file_name=f"Business_Plan_{selected_scenario.replace(' ', '_')}.docx",
                        mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                        key=f"download_business_plan_docx_{selected_scenario}",
                    )
            with export_cols[1]:
                try:
                    pdf_bytes = _generate_business_plan_pdf_bytes(
                        export_title,
                        plan_markdown,
                        investor_recommendations,
                        citation_entries,
                    )
                except RuntimeError as exc:
                    st.warning(str(exc))
                else:
                    st.download_button(
                        "Download PDF (.pdf)",
                        data=pdf_bytes,
                        file_name=f"Business_Plan_{selected_scenario.replace(' ', '_')}.pdf",
                        mime="application/pdf",
                        key=f"download_business_plan_pdf_{selected_scenario}",
                    )
            with export_cols[2]:
                business_plan_tables: Dict[str, pd.DataFrame] = {
                    "Assumptions Schedule": assumption_schedule_df,
                    "Project Input Values": pd.DataFrame([asdict(assumptions)]),
                    "Production Cycles": cycles_df,
                    "Annual Summary": annual_df,
                    "Cash Flow Projections": cashflow_df,
                    "Valuation Analysis": valuation_df,
                    "Income Statement": income_df,
                    "Statement of Financial Position": balance_df,
                    "Cash Flow Statement": cash_statement_df,
                    "Debt Schedule": loan_df,
                    "Revenue Summary by Year": annual_totals_df,
                    "Revenue Summary by Category": summary_by_category,
                    "Advanced Metrics": metrics_df,
                    "DSCR Schedule": dscr_df,
                    "Trend Analysis": trend_df,
                    "Interest Coverage": coverage_df,
                    "Leverage Metrics": leverage_df,
                    "What If Analysis": what_if_df,
                    "Break Even Analysis": break_even_df,
                    "Scenario Planning": scenario_df,
                    "Forecast Projections": forecast_df,
                    "Time Series Projections": time_series_df,
                    "Risk and Anomaly Log": risk_df,
                    "ML Methods": ml_methods_df,
                    "Monte Carlo Summary": monte_carlo_summary_df,
                    "Monte Carlo Samples": monte_carlo_samples_df,
                }
                for category, rows in revenue_schedules.items():
                    safe_key = f"Revenue Schedule {category}"
                    business_plan_tables[safe_key] = pd.DataFrame(rows)
                business_plan_tables = {
                    key: value
                    for key, value in business_plan_tables.items()
                    if isinstance(value, pd.DataFrame) and not value.empty
                }

                revenue_chart = _build_revenue_stack_chart(summary_by_category)
                income_mix_chart = _build_income_composition_chart(income_df)
                leverage_chart = _combined_leverage_chart(dscr_df, coverage_df, leverage_df)

                mc_chart: Optional[alt.Chart] = None
                if not monte_carlo_samples_df.empty and {"iteration", "npv"}.issubset(
                    monte_carlo_samples_df.columns
                ):
                    mc_chart = (
                        alt.Chart(monte_carlo_samples_df)
                        .mark_bar(opacity=0.8)
                        .encode(
                            x=alt.X("npv:Q", bin=alt.Bin(maxbins=40), title="NPV (USD)"),
                            y=alt.Y("count():Q", title="Frequency"),
                            tooltip=[alt.Tooltip("count():Q", title="Samples")],
                        )
                    )

                chart_payloads: Dict[str, Dict[str, Any]] = {}
                if revenue_chart is not None:
                    chart_payloads["Revenue Mix and Growth"] = {
                        "data": summary_by_category,
                        "spec": revenue_chart.to_dict(),
                    }
                if income_mix_chart is not None:
                    chart_payloads["Profitability Profile"] = {
                        "data": income_df,
                        "spec": income_mix_chart.to_dict(),
                    }
                if leverage_chart is not None:
                    leverage_chart_data = pd.concat(
                        [df for df in [dscr_df, coverage_df, leverage_df] if not df.empty],
                        axis=1,
                    )
                    chart_payloads["Leverage and Resilience"] = {
                        "data": leverage_chart_data,
                        "spec": leverage_chart.to_dict(),
                    }
                if mc_chart is not None:
                    chart_payloads["NPV Risk Distribution"] = {
                        "data": monte_carlo_samples_df,
                        "spec": mc_chart.to_dict(),
                    }

                try:
                    plan_excel_bytes = _generate_business_plan_excel_bytes(
                        selected_scenario,
                        plan_markdown,
                        investor_recommendations,
                        scorecard_df,
                        confidence_df,
                        citation_entries,
                        financial_tables=business_plan_tables,
                        chart_payloads=chart_payloads,
                    )
                except RuntimeError as exc:
                    if "Missing Excel writer dependency" in str(exc):
                        plan_csv_zip_bytes = _generate_business_plan_csv_zip_bytes(
                            selected_scenario,
                            plan_markdown,
                            investor_recommendations,
                            scorecard_df,
                            confidence_df,
                            citation_entries,
                            financial_tables=business_plan_tables,
                            chart_payloads=chart_payloads,
                        )
                        st.warning(
                            "Excel writer dependencies are unavailable. Downloading CSV ZIP export instead."
                        )
                        st.download_button(
                            "Download Business Plan CSV ZIP (.zip)",
                            data=plan_csv_zip_bytes,
                            file_name=f"Business_Plan_{selected_scenario.replace(' ', '_')}_CSV_Export.zip",
                            mime="application/zip",
                            key=f"download_business_plan_csv_zip_{selected_scenario}",
                        )
                    else:
                        st.warning(str(exc))
                else:
                    st.download_button(
                        "Download Excel (.xlsx)",
                        data=plan_excel_bytes,
                        file_name=f"Business_Plan_{selected_scenario.replace(' ', '_')}.xlsx",
                        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                        key=f"download_business_plan_excel_{selected_scenario}",
                    )

            summary_by_category = pd.DataFrame(revenue_summary.get("by_category", []))
            revenue_chart = _build_revenue_stack_chart(summary_by_category)
            if revenue_chart is not None:
                st.markdown("#### Revenue mix and growth outlook")
                st.altair_chart(revenue_chart, use_container_width=True)

            income_mix_chart = _build_income_composition_chart(income_df)
            if income_mix_chart is not None:
                st.markdown("#### Profitability profile")
                st.altair_chart(income_mix_chart, use_container_width=True)

            leverage_chart = _combined_leverage_chart(dscr_df, coverage_df, leverage_df)
            if leverage_chart is not None:
                st.markdown("#### Leverage and debt-service resilience")
                st.altair_chart(leverage_chart, use_container_width=True)

            if mc_chart is not None:
                st.markdown("#### NPV risk distribution")
                st.altair_chart(mc_chart, use_container_width=True)

        unified_ai_state["outputs"]["chat_messages"] = len(chat_history)
        unified_ai_state["outputs"]["business_plan_enabled"] = bool(
            st.session_state.get(plan_state_key, False)
        )
        unified_ai_state["outputs"]["timestamp_utc"] = datetime.now(
            timezone.utc
        ).isoformat()
        payload["ai_unified_system"] = unified_ai_state
        st.session_state["scenario_store"][selected_scenario] = payload

    analytics_namespace = "advanced_schedule_state"

    with analytics_tab:
        if not metrics_df.empty:
            st.subheader("Advanced metrics")
            metrics_df = _render_analytics_schedule(
                "Advanced metrics",
                "advanced_metrics",
                metrics_df,
                selected_scenario,
                namespace=analytics_namespace,
            )
            advanced["metrics"] = metrics_df.replace({pd.NA: None}).to_dict("records")

        monte_carlo_p5 = None
        monte_carlo_p50 = None
        monte_carlo_p95 = None
        if not monte_carlo_summary_df.empty:
            summary_row = monte_carlo_summary_df.iloc[0]
            monte_carlo_p5 = _to_float(summary_row.get("p5_npv"))
            monte_carlo_p95 = _to_float(summary_row.get("p95_npv"))
        if not monte_carlo_samples_df.empty and "npv" in monte_carlo_samples_df.columns:
            npv_series = pd.to_numeric(monte_carlo_samples_df["npv"], errors="coerce").dropna()
            if not npv_series.empty:
                monte_carlo_p50 = float(npv_series.quantile(0.5))

        avg_break_even_price = None
        if not break_even_df.empty and "Break-even price" in break_even_df.columns:
            break_even_series = pd.to_numeric(
                break_even_df["Break-even price"], errors="coerce"
            ).dropna()
            if not break_even_series.empty:
                avg_break_even_price = float(break_even_series.mean())

        annual_net_income = None
        if not income_df.empty and "net_income" in income_df.columns:
            income_series = pd.to_numeric(income_df["net_income"], errors="coerce").dropna()
            if not income_series.empty:
                annual_net_income = float(income_series.iloc[0])

        payback_period = None
        if not metrics_df.empty and {"metric", "value"}.issubset(metrics_df.columns):
            payback_row = metrics_df.loc[
                metrics_df["metric"] == "Payback period (years)", "value"
            ]
            if not payback_row.empty:
                payback_period = _to_float(payback_row.iloc[0])
        st.markdown("### Simulation diagnostics")
        st.markdown(
            "\n".join(
                [
                    f"Monte Carlo NPV P5: { _format_currency(monte_carlo_p5) }",
                    f"Monte Carlo NPV P50: { _format_currency(monte_carlo_p50) }",
                    f"Monte Carlo NPV P95: { _format_currency(monte_carlo_p95) }",
                    f"Average break-even price per kg: { _format_currency(avg_break_even_price) }",
                    f"Annual net income (current run): { _format_currency(annual_net_income) }",
                    f"Payback period: {payback_period:.2f} years" if payback_period is not None else "Payback period: N/A years",
                ]
            )
        )

        st.subheader("Simulation builder")
        custom_defaults = copy.deepcopy(custom_definition_defaults)
        custom_editor_df = _render_schedule_editor(
            "Custom scenario definitions",
            "custom_simulations",
            custom_defaults,
            selected_scenario,
            namespace="custom_simulation_state",
            row_defaults={
                "Scenario": "New scenario",
                "Description": "",
                "Parameter": "",
                "Change type": "percent",
                "Change value": 0.0,
            },
            allow_yearly_increment=False,
        )
        custom_rows = (
            custom_editor_df.replace({pd.NA: None}).to_dict("records")
            if not custom_editor_df.empty
            else []
        )
        simulation_inputs = custom_rows or copy.deepcopy(custom_definition_defaults)
        stored_payload = custom_simulations_data or {}
        stored_definitions = stored_payload.get("definitions", [])
        run_custom_button = st.button(
            "Run custom simulations",
            key=f"run_custom_simulations_{selected_scenario}",
            help="Recalculate the scenarios above with the current definitions.",
        )

        simulation_payload = stored_payload
        if run_custom_button:
            with st.spinner("Running custom simulations..."):
                simulation_payload = run_custom_simulations(
                    model.assumptions,
                    base_metrics,
                    simulation_inputs,
                )
            custom_simulations_data = simulation_payload
            advanced["custom_simulations"] = simulation_payload
            advanced["custom_simulation_definitions"] = copy.deepcopy(simulation_inputs)
            results["advanced_analytics"]["custom_simulations"] = copy.deepcopy(
                simulation_payload
            )
            results["advanced_analytics"][
                "custom_simulation_definitions"
            ] = copy.deepcopy(simulation_inputs)
            payload_cache = st.session_state.get("scenario_payloads", {})
            if selected_scenario in payload_cache:
                payload_cache[selected_scenario]["results"]["advanced_analytics"][
                    "custom_simulations"
                ] = copy.deepcopy(simulation_payload)
                payload_cache[selected_scenario]["results"]["advanced_analytics"][
                    "custom_simulation_definitions"
                ] = copy.deepcopy(simulation_inputs)
        else:
            if not _records_match(simulation_inputs, stored_definitions):
                st.info(
                    "Definitions changed. Click **Run custom simulations** to refresh the "
                    "results below."
                )

        custom_state_store = st.session_state.setdefault("custom_simulation_state", {})
        scenario_store_custom = custom_state_store.setdefault(selected_scenario, {})
        schedule_state = scenario_store_custom.setdefault("custom_simulations", {})
        schedule_state["data"] = copy.deepcopy(simulation_inputs)
        schedule_state.setdefault("default", copy.deepcopy(custom_defaults))

        custom_results_df = pd.DataFrame(simulation_payload.get("results", []))
        if custom_rows:
            valid_count = len(simulation_payload.get("definitions", []))
            configured = sum(1 for row in custom_rows if row.get("Scenario"))
            if run_custom_button and valid_count < configured:
                st.warning(
                    "Some rows were skipped because the parameter name or change value "
                    "could not be applied. Check the entries above and try again."
                )
        invalid_rows = simulation_payload.get("invalid", [])
        if invalid_rows:
            st.info(
                "Rows with validation issues are retained above but excluded from the "
                "simulation results."
            )
            invalid_df = pd.DataFrame(invalid_rows)
            st.dataframe(
                invalid_df,
                use_container_width=True,
                hide_index=True,
            )
        if not custom_results_df.empty:
            st.dataframe(
                custom_results_df,
                use_container_width=True,
                hide_index=True,
            )
            delta_df = pd.DataFrame(simulation_payload.get("delta_summary", []))
            if not delta_df.empty:
                st.caption("NPV delta by custom scenario")
                try:
                    chart_series = delta_df.set_index("Scenario")["NPV Δ"].astype(float)
                    st.bar_chart(chart_series)
                except KeyError:
                    pass

        st.subheader("What-if analysis")
        run_what_if = st.button(
            "Run what-if analysis",
            key=f"run_what_if_{selected_scenario}",
            help="Recompute the what-if scenarios using the latest assumptions.",
        )
        if run_what_if:
            with st.spinner("Running what-if scenarios..."):
                what_if_rows = perform_what_if_analysis(
                    model.assumptions,
                    base_metrics,
                )
            what_if_df = pd.DataFrame(what_if_rows)
            advanced["what_if"] = what_if_rows
            results["advanced_analytics"]["what_if"] = copy.deepcopy(what_if_rows)
            payload_cache = st.session_state.get("scenario_payloads", {})
            if selected_scenario in payload_cache:
                payload_cache[selected_scenario]["results"]["advanced_analytics"][
                    "what_if"
                ] = copy.deepcopy(what_if_rows)
        if what_if_df.empty:
            st.info(
                "No what-if scenarios available yet. Click **Run what-if analysis** to "
                "generate the table."
            )
        else:
            what_if_df = _render_analytics_schedule(
                "What-if scenarios",
                "what_if",
                what_if_df,
                selected_scenario,
                namespace=analytics_namespace,
            )
            advanced["what_if"] = what_if_df.replace({pd.NA: None}).to_dict("records")

        st.subheader("Scenario planning")
        run_scenario_planning = st.button(
            "Run scenario planning",
            key=f"run_scenario_planning_{selected_scenario}",
            help="Recalculate Downside, Upside, and Expansion scenarios.",
        )
        if run_scenario_planning:
            with st.spinner("Evaluating planning scenarios..."):
                scenario_rows = scenario_planning(
                    model.assumptions,
                    base_metrics,
                    getattr(annual_summary_obj, "revenue", 0.0),
                )
            scenario_df = pd.DataFrame(scenario_rows)
            advanced["scenario_planning"] = scenario_rows
            results["advanced_analytics"]["scenario_planning"] = copy.deepcopy(
                scenario_rows
            )
            payload_cache = st.session_state.get("scenario_payloads", {})
            if selected_scenario in payload_cache:
                payload_cache[selected_scenario]["results"]["advanced_analytics"][
                    "scenario_planning"
                ] = copy.deepcopy(scenario_rows)
        if scenario_df.empty:
            st.info(
                "No scenario planning results yet. Click **Run scenario planning** to "
                "evaluate alternative cases."
            )
        else:
            scenario_df = _render_analytics_schedule(
                "Scenario planning",
                "scenario_planning",
                scenario_df,
                selected_scenario,
                namespace=analytics_namespace,
            )
            advanced["scenario_planning"] = (
                scenario_df.replace({pd.NA: None}).to_dict("records")
            )

        tornado_chart = _tornado_chart(what_if_df, scenario_df)
        if tornado_chart is not None:
            st.markdown("#### Sensitivity tornado (NPV deltas)")
            st.altair_chart(tornado_chart, use_container_width=True)

        monte_carlo_settings = (
            monte_carlo.get("settings", {}) if isinstance(monte_carlo, dict) else {}
        )
        default_iterations = int(
            monte_carlo_settings.get("iterations")
            or (
                int(monte_carlo_summary_df.iloc[0].get("iterations"))
                if not monte_carlo_summary_df.empty
                and pd.notna(monte_carlo_summary_df.iloc[0].get("iterations"))
                else len(monte_carlo_samples_df) or 200
            )
        )
        default_iterations = min(max(default_iterations, 1), 10000)

        distribution_defaults = monte_carlo_settings.get("distributions")
        if not distribution_defaults:
            distribution_defaults = copy.deepcopy(DEFAULT_MONTE_CARLO_DISTRIBUTIONS)
        else:
            distribution_defaults = copy.deepcopy(distribution_defaults)
        distributions_df = pd.DataFrame(distribution_defaults)

        st.subheader("Monte Carlo simulation")
        st.caption("Adjust parameter distributions before running simulations.")
        distributions_df = _render_analytics_schedule(
            "Monte Carlo distributions",
            "monte_carlo_distributions",
            distributions_df,
            selected_scenario,
            namespace=analytics_namespace,
            iterable_columns=["bounds"],
        )
        distribution_records = (
            distributions_df.replace({pd.NA: None}).to_dict("records")
            if not distributions_df.empty
            else []
        )
        if not distribution_records:
            distribution_records = copy.deepcopy(DEFAULT_MONTE_CARLO_DISTRIBUTIONS)
        monte_carlo_settings["distributions"] = copy.deepcopy(distribution_records)

        iterations_key = f"monte_carlo_iterations_{selected_scenario}"
        if iterations_key not in st.session_state:
            st.session_state[iterations_key] = default_iterations

        iterations_value = st.number_input(
            "Iterations",
            min_value=1,
            max_value=10000,
            value=st.session_state[iterations_key],
            step=50,
            key=iterations_key,
            help="Choose how many Monte Carlo iterations to run (1-10,000).",
        )

        run_button = st.button(
            "Run Monte Carlo",
            key=f"run_monte_carlo_{selected_scenario}",
        )

        if run_button:
            iterations = int(iterations_value)
            with st.spinner("Running Monte Carlo simulation..."):
                updated_mc = run_monte_carlo_analysis(
                    model.assumptions,
                    iterations=iterations,
                    distributions=distribution_records,
                )
            monte_carlo = updated_mc
            advanced["monte_carlo"] = updated_mc
            results["advanced_analytics"]["monte_carlo"] = copy.deepcopy(updated_mc)
            advanced["monte_carlo"].setdefault("settings", {})["distributions"] = (
                copy.deepcopy(distribution_records)
            )
            payload_cache = st.session_state.get("scenario_payloads", {})
            if selected_scenario in payload_cache:
                payload_cache[selected_scenario]["results"]["advanced_analytics"][
                    "monte_carlo"
                ] = copy.deepcopy(updated_mc)
            monte_carlo_summary_df = pd.DataFrame([updated_mc.get("summary", {})])
            monte_carlo_samples_df = pd.DataFrame(updated_mc.get("samples", []))
            monte_carlo_settings = updated_mc.get("settings", {})

        if not monte_carlo_summary_df.empty:
            st.markdown("#### Summary")
            st.dataframe(
                monte_carlo_summary_df,
                use_container_width=True,
                hide_index=True,
            )

        if not monte_carlo_samples_df.empty:
            st.caption("Monte Carlo outcome distributions")
            for label, chart in _monte_carlo_distribution_charts(
                monte_carlo_samples_df
            ):
                st.markdown(f"**{label} distribution**")
                st.altair_chart(chart, use_container_width=True)

        st.subheader("Break-even analysis by product")
        run_break_even = st.button(
            "Run break-even analysis",
            key=f"run_break_even_{selected_scenario}",
            help="Rebuild the break-even schedule using current revenue inputs.",
        )
        if run_break_even:
            with st.spinner("Calculating break-even metrics..."):
                break_even_rows = break_even_analysis(
                    annual_summary_obj,
                    revenue_summary,
                    revenue_schedules,
                )
            break_even_df = pd.DataFrame(break_even_rows)
            advanced["break_even"] = break_even_rows
            results["advanced_analytics"]["break_even"] = copy.deepcopy(break_even_rows)
            payload_cache = st.session_state.get("scenario_payloads", {})
            if selected_scenario in payload_cache:
                payload_cache[selected_scenario]["results"]["advanced_analytics"][
                    "break_even"
                ] = copy.deepcopy(break_even_rows)
        if break_even_df.empty:
            st.info(
                "No break-even schedule is available. Click **Run break-even analysis** to "
                "generate the table."
            )
        else:
            break_even_df = _render_analytics_schedule(
                "Break-even analysis",
                "break_even",
                break_even_df,
                selected_scenario,
                namespace=analytics_namespace,
            )
            advanced["break_even"] = (
                break_even_df.replace({pd.NA: None}).to_dict("records")
            )
            heatmap_chart, units_chart = _break_even_heatmap_charts(break_even_df)
            if heatmap_chart is not None:
                st.altair_chart(heatmap_chart, use_container_width=True)
            if units_chart is not None:
                st.altair_chart(units_chart, use_container_width=True)

        st.subheader("Goal seek (target NPV)")
        run_goal_seek = st.button(
            "Run goal seek",
            key=f"run_goal_seek_{selected_scenario}",
            help="Solve for the live price per kg that delivers zero NPV.",
        )
        if run_goal_seek:
            with st.spinner("Running goal seek..."):
                goal_seek = goal_seek_live_price(model.assumptions)
            advanced["goal_seek"] = goal_seek
            results["advanced_analytics"]["goal_seek"] = copy.deepcopy(goal_seek)
            payload_cache = st.session_state.get("scenario_payloads", {})
            if selected_scenario in payload_cache:
                payload_cache[selected_scenario]["results"]["advanced_analytics"][
                    "goal_seek"
                ] = copy.deepcopy(goal_seek)
        goal_df = pd.DataFrame([goal_seek]) if goal_seek else pd.DataFrame()
        if goal_df.empty:
            st.info(
                "No goal seek result yet. Click **Run goal seek** to calculate the target "
                "price."
            )
        else:
            goal_df = _render_analytics_schedule(
                "Goal seek",
                "goal_seek",
                goal_df,
                selected_scenario,
                namespace=analytics_namespace,
            )
            goal_records = goal_df.replace({pd.NA: None}).to_dict("records")
            advanced["goal_seek"] = goal_records[0] if goal_records else {}

        if not dscr_df.empty:
            st.subheader("Debt service coverage ratio")
            dscr_df = _render_analytics_schedule(
                "DSCR summary",
                "dscr",
                dscr_df,
                selected_scenario,
                namespace=analytics_namespace,
            )
            advanced["dscr"] = dscr_df.replace({pd.NA: None}).to_dict("records")

        if not returns_df.empty:
            st.subheader("Return diagnostics")
            returns_df = _render_analytics_schedule(
                "Return diagnostics",
                "return_diagnostics",
                returns_df,
                selected_scenario,
                namespace=analytics_namespace,
            )
            advanced["returns"] = returns_df.replace({pd.NA: None}).to_dict("records")
            try:
                returns_chart = returns_df.set_index("year")
                returns_chart = returns_chart.apply(pd.to_numeric, errors="coerce")
                subset_cols = [
                    col
                    for col in [
                        "return_on_assets",
                        "return_on_equity",
                        "return_on_invested_capital",
                    ]
                    if col in returns_chart.columns
                ]
                if subset_cols:
                    st.line_chart(returns_chart[subset_cols].dropna(how="all"))
            except KeyError:
                pass

        if not coverage_df.empty:
            st.subheader("Coverage & resilience")
            coverage_df = _render_analytics_schedule(
                "Coverage & resilience",
                "coverage",
                coverage_df,
                selected_scenario,
                namespace=analytics_namespace,
            )
            advanced["coverage"] = (
                coverage_df.replace({pd.NA: None}).to_dict("records")
            )

        if not leverage_df.empty:
            st.subheader("Leverage profile")
            leverage_df = _render_analytics_schedule(
                "Leverage profile",
                "leverage",
                leverage_df,
                selected_scenario,
                namespace=analytics_namespace,
            )
            advanced["leverage"] = (
                leverage_df.replace({pd.NA: None}).to_dict("records")
            )

        combined_leverage_chart = _combined_leverage_chart(
            dscr_df, coverage_df, leverage_df
        )
        if combined_leverage_chart is not None:
            st.markdown("#### Coverage vs leverage overview")
            st.altair_chart(combined_leverage_chart, use_container_width=True)

        if not trend_df.empty:
            st.subheader("Performance trends")
            trend_df = _render_analytics_schedule(
                "Performance trends",
                "performance_trends",
                trend_df,
                selected_scenario,
                namespace=analytics_namespace,
            )
            advanced["trend"] = trend_df.replace({pd.NA: None}).to_dict("records")
            try:
                trend_chart = trend_df.set_index("year")
                trend_chart = trend_chart.apply(pd.to_numeric, errors="coerce")
                subset_cols = [
                    col
                    for col in [
                        "revenue",
                        "ebitda",
                        "net_income",
                        "free_cash_flow",
                    ]
                    if col in trend_chart.columns
                ]
                if subset_cols:
                    st.line_chart(trend_chart[subset_cols].dropna(how="all"))
            except KeyError:
                pass

        st.subheader("Predictive analytics")
        auto_run_predictive = bool(ai_settings.get("auto_run_predictive", False))
        run_predictive_clicked = st.button(
            "Run predictive analytics",
            key=f"run_predictive_{selected_scenario}",
            help="Refresh automated forecasts, time-series projections, and risk checks.",
        )
        if auto_run_predictive:
            st.caption("Auto-run predictive analytics is enabled for this scenario.")
        run_predictive = run_predictive_clicked or auto_run_predictive
        if run_predictive:
            with st.spinner("Recomputing predictive analytics..."):
                predictive = build_predictive_analytics(
                    results["cashflows"],
                    financials["income_statement"],
                )
            advanced["predictive"] = predictive
            results["advanced_analytics"]["predictive"] = copy.deepcopy(predictive)
            payload_cache = st.session_state.get("scenario_payloads", {})
            if selected_scenario in payload_cache:
                payload_cache[selected_scenario]["results"]["advanced_analytics"][
                    "predictive"
                ] = copy.deepcopy(predictive)
        forecast_df = pd.DataFrame(predictive.get("automated_forecast", []))
        time_series_df = pd.DataFrame(
            predictive.get("time_series", {}).get("forecast", [])
        )
        risk_metadata = predictive.get("risk_anomalies", {})
        risk_df = pd.DataFrame(risk_metadata.get("observations", []))
        ml_methods_df = pd.DataFrame(predictive.get("ml_methods", []))

        if (
            forecast_df.empty
            and time_series_df.empty
            and risk_df.empty
            and ml_methods_df.empty
        ):
            st.info(
                "No predictive analytics available. Click **Run predictive analytics** to "
                "generate forecasts and anomaly checks."
            )

        if not forecast_df.empty:
            st.markdown("##### Automated forecasting")
            forecast_df = _render_analytics_schedule(
                "Automated forecast",
                "automated_forecast",
                forecast_df,
                selected_scenario,
                namespace=analytics_namespace,
            )
            predictive.setdefault("automated_forecast", [])
            predictive["automated_forecast"] = (
                forecast_df.replace({pd.NA: None}).to_dict("records")
            )
            try:
                forecast_chart = forecast_df.set_index("Year")
                forecast_chart = forecast_chart.apply(pd.to_numeric, errors="coerce")
                subset_cols = [
                    col
                    for col in ["Revenue forecast", "EBITDA forecast"]
                    if col in forecast_chart.columns
                ]
                if subset_cols:
                    st.line_chart(forecast_chart[subset_cols].dropna(how="all"))
            except KeyError:
                pass

        if not time_series_df.empty:
            st.markdown("##### Time series (AR(1)) outlook")
            time_series_df = _render_analytics_schedule(
                "Time series outlook",
                "time_series",
                time_series_df,
                selected_scenario,
                namespace=analytics_namespace,
            )
            predictive.setdefault("time_series", {}).setdefault("forecast", [])
            predictive["time_series"]["forecast"] = (
                time_series_df.replace({pd.NA: None}).to_dict("records")
            )
            try:
                time_series_chart = time_series_df.set_index("Year")
                time_series_chart = time_series_chart.apply(pd.to_numeric, errors="coerce")
                st.line_chart(time_series_chart.dropna(how="all"))
            except KeyError:
                pass

        if not risk_df.empty:
            st.markdown("##### Risk & anomaly detection")
            risk_df = _render_analytics_schedule(
                "Risk & anomalies",
                "risk_anomalies",
                risk_df,
                selected_scenario,
                namespace=analytics_namespace,
            )
            predictive.setdefault("risk_anomalies", {})
            predictive["risk_anomalies"]["observations"] = (
                risk_df.replace({pd.NA: None}).to_dict("records")
            )
            st.caption(
                f"Mean growth: {risk_metadata.get('mean_growth', float('nan')):.2%} — "
                f"Std dev: {risk_metadata.get('std_growth', float('nan')):.2%}"
            )

        anomaly_points = _build_anomaly_points(risk_df, trend_df)
        overlay_chart = _forecast_overlay_chart(
            trend_df, forecast_df, time_series_df, anomaly_points
        )
        if overlay_chart is not None:
            st.markdown("#### Revenue outlook (historical vs forecasts)")
            st.altair_chart(overlay_chart, use_container_width=True)

        if not ml_methods_df.empty:
            st.markdown("##### ML method diagnostics")
            ml_methods_df = _render_analytics_schedule(
                "ML methods",
                "ml_methods",
                ml_methods_df,
                selected_scenario,
                namespace=analytics_namespace,
            )
            predictive["ml_methods"] = (
                ml_methods_df.replace({pd.NA: None}).to_dict("records")
            )

        advanced["monte_carlo"] = monte_carlo
        advanced["predictive"] = predictive

    st.session_state.input_snapshot = copy.deepcopy(payload)
    scenario_store[selected_scenario] = payload
    scenario_payloads = st.session_state.get("scenario_payloads")
    if scenario_payloads and selected_scenario in scenario_payloads:
        scenario_payloads[selected_scenario]["snapshot"] = copy.deepcopy(payload)
        scenario_payloads[selected_scenario]["results"] = results

    st.markdown("---")
    st.caption("Use the Input Landing Page above to adjust the operating model and financing structure.")


if __name__ == "__main__":
    main()
